import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────────
BG         = RGBColor(0x0F, 0x17, 0x2A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PL         = RGBColor(0xA7, 0x8B, 0xFA)  # purple light
PM         = RGBColor(0x8B, 0x5C, 0xF6)  # purple mid
PD         = RGBColor(0x6D, 0x28, 0xD9)  # purple deep
SL         = RGBColor(0xE2, 0xE8, 0xF0)  # slate light
SM         = RGBColor(0x33, 0x41, 0x55)  # slate mid
SD         = RGBColor(0x1E, 0x29, 0x3B)  # slate dark (cards)
GREEN      = RGBColor(0x34, 0xD3, 0x99)
YELLOW     = RGBColor(0xFB, 0xBF, 0x24)

SLIDE_W = 12192000
SLIDE_H  = 6858000

def cm(v):    return int(v * 360000)
def pt(v):    return Pt(v)

# ── Primitives ───────────────────────────────────────────────
def bg(slide, color=BG):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = color
    else:
        s.line.fill.background()
    return s

def tb(slide, x, y, w, h, text, size, color,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.name = 'Inter'; f.size = pt(size)
    f.color.rgb = color; f.bold = bold; f.italic = italic
    return box

def title_bar(slide, text, size=34):
    tb(slide, 2.5, 1.3, 27.9, 1.9, text, size, WHITE)

def purple_bar(slide, top=3.3, h=2.0):
    rect(slide, 2.0, top, 0.2, h, PD)

def footer(slide, text, color=SM, size=12):
    tb(slide, 2.5, 17.0, 27.9, 1.0, text, size, color)

def multiline_tb(slide, x, y, w, h, lines):
    """lines: list of (text, size, color, bold)"""
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for (text, size, color, bold) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if text == '':
            r = p.add_run(); r.text = ''
            r.font.name = 'Inter'; r.font.size = pt(size)
            continue
        r = p.add_run(); r.text = text
        f = r.font; f.name = 'Inter'; f.size = pt(size)
        f.color.rgb = color; f.bold = bold
    return box

def col_card(slide, x, y_top, w, h, period, title, color, items, item_size=13):
    rect(slide, x, y_top, w, h, SD)
    rect(slide, x, y_top, w, 0.15, color)
    tb(slide, x+0.5, y_top+0.3, w-0.6, 0.7, period, 11, color)
    tb(slide, x+0.5, y_top+1.1, w-0.6, 1.5, title, 18, WHITE)
    lines = []
    for i, item in enumerate(items):
        if i > 0:
            lines.append(('', 4, SL, False))
        lines.append((f'\u2192  {item}', item_size, SL, False))
    multiline_tb(slide, x+0.5, y_top+3.0, w-0.6, h-3.5, lines)

# ── Build ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BL = prs.slide_layouts[6]  # blank

# ================================================================
# SLIDE 1 — Cover
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
tb(s, 2.5, 3.2, 25.4, 1.8, 'ELIOS SCOGLIO', 46, WHITE)
tb(s, 2.5, 5.2, 27.9, 0.9, 'Software & Architecture Manager  |  Fractional CTO', 22, PL)
rect(s, 2.5, 6.5, 8.1, 0.1, PD)
# Subtitle line
box = s.shapes.add_textbox(cm(2.5), cm(7.2), cm(27.9), cm(1.2))
tf = box.text_frame
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = 'Secondo colloquio — '
r1.font.name='Inter'; r1.font.size=Pt(17); r1.font.color.rgb=SL
r2 = p.add_run(); r2.text = 'Veralab'
r2.font.name='Inter'; r2.font.size=Pt(17); r2.font.color.rgb=PL; r2.font.bold=True

tb(s, 2.5, 9.0, 27.9, 0.9,
   '"Operativo nella transizione. Strategico nel regime."',
   15, SL, italic=True)
tb(s, 2.5, 17.5, 27.9, 0.7,
   '108 Vision \u2014 Costruiamo la direzione, non solo il codice.',
   13, SM)

# ================================================================
# SLIDE 2 — Perché questa conversazione
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Perch\u00e9 questa conversazione')
purple_bar(s)

lines = [
    ('Veralab non \u00e8 pi\u00f9 soltanto un e-commerce.', 16, SL, True),
    ('', 6, SL, False),
    ('\u00c8 un ecosistema che deve collegare:', 15, SL, False),
    ('', 6, SL, False),
    ('\u2192  commercio digitale e punti vendita fisici', 14, SL, False),
    ('\u2192  ERP e logistica', 14, SL, False),
    ('\u2192  CRM e loyalty', 14, SL, False),
    ('\u2192  marketing, dati cliente, nuovi brand', 14, SL, False),
    ('\u2192  crescita nazionale e internazionale', 14, SL, False),
    ('', 8, SL, False),
    ("L'azienda aumenta il retail fisico, investe in CRM & Data e prepara una crescita guidata dai dati.", 14, SM, False),
    ('', 8, SL, False),
    ('In questa fase la tecnologia deve diventare il sistema operativo della crescita.', 16, PL, True),
]
multiline_tb(s, 3.0, 4.3, 27.9, 13.0, lines)

# ================================================================
# SLIDE 3 — La mia proposta
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'La mia proposta')
purple_bar(s)

lines = [
    ('Non una consulenza esterna che osserva e produce documenti.', 16, SL, True),
    ('Un Fractional CTO operativo, capace di:', 15, SL, False),
    ('', 6, SL, False),
    ('\u2192  prendere ownership delle priorit\u00e0 tecniche', 15, SL, False),
    ('\u2192  collegare strategia, tecnologia e business', 15, SL, False),
    ('\u2192  intervenire direttamente sui problemi critici', 15, SL, False),
    ('\u2192  supportare il responsabile tecnico e il team', 15, SL, False),
    ('\u2192  coordinare fornitori e piattaforme', 15, SL, False),
    ('\u2192  costruire processi decisionali pi\u00f9 chiari', 15, SL, False),
    ("\u2192  rendere progressivamente autonoma l'organizzazione", 15, SL, False),
    ('', 8, SL, False),
    ('Nella fase iniziale il coinvolgimento sar\u00e0 pi\u00f9 intenso. Con il consolidamento il presidio verr\u00e0 progressivamente ridotto e concentrato sulle decisioni strategiche.', 13, SM, False),
    ('', 6, SL, False),
    ('"Il risultato non \u00e8 rendermi indispensabile. \u00c8 rendere l\'organizzazione pi\u00f9 capace anche quando non sono presente."', 15, PL, False),
]
multiline_tb(s, 3.0, 4.3, 27.9, 13.5, lines)

# ================================================================
# SLIDE 4 — Valore dal primo giorno
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Valore dal primo giorno')
purple_bar(s)
tb(s, 3.0, 4.3, 27.9, 0.8, 'Comprendo e intervengo contemporaneamente \u2014 dal giorno 1.', 15, SL)

lines_l = [
    ('INTERVENTI IMMEDIATI', 12, PL, True),
    ('', 4, SL, False),
    ('\u2192  Rendere visibili errori e anomalie', 13, SL, False),
    ('\u2192  Chiarire chi decide e chi \u00e8 responsabile', 13, SL, False),
    ('\u2192  Eliminare blocchi decisionali', 13, SL, False),
    ('\u2192  Ridurre i rischi dei rilasci', 13, SL, False),
    ('\u2192  Verificare affidabilit\u00e0 delle integrazioni', 13, SL, False),
    ('\u2192  Introdurre gestione incidenti minima', 13, SL, False),
    ('\u2192  Eliminare attivit\u00e0 manuali ripetitive', 13, SL, False),
    ('\u2192  Ordinare backlog per impatto sul business', 13, SL, False),
]
lines_r = [
    ('DIFFERENZE CHE CONTANO', 12, GREEN, True),
    ('', 4, SL, False),
    ('\u2713  Un alert prima della segnalazione cliente', 13, SL, False),
    ('\u2713  Un rilascio reversibile', 13, SL, False),
    ('\u2713  Un processo con un responsabile chiaro', 13, SL, False),
    ('\u2713  Un\'integrazione che non perde eventi', 13, SL, False),
    ('\u2713  Un dato cliente affidabile', 13, SL, False),
    ('\u2713  Una decisione documentata', 13, SL, False),
    ('\u2713  Una campagna senza emergenza tecnica', 13, SL, False),
    ('\u2713  Un problema ricorrente che smette di ripetersi', 13, SL, False),
]
multiline_tb(s, 2.5, 5.5, 14.0, 12.0, lines_l)
multiline_tb(s, 17.5, 5.5, 13.5, 12.0, lines_r)

# ================================================================
# SLIDE 5 — I 5 pilastri
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'I cinque pilastri del mio lavoro')
purple_bar(s)

pillars = [
    ('Strategia',
     'Trasformare gli obiettivi aziendali in priorit\u00e0 tecnologiche comprensibili, misurabili e sostenibili.'),
    ('Architettura',
     'Costruire sistemi capaci di sostenere crescita, nuovi canali e processi senza moltiplicare fragilit\u00e0.'),
    ('Team',
     'Creare ownership, chiarezza e capacit\u00e0 decisionale. Leva per il responsabile tecnico, non livello che lo rallenta.'),
    ('Governance',
     'Rendere visibili decisioni, costi, rischi e risultati. Mai sorprese. Sempre opzioni. Trade-off espliciti.'),
    ('AI & Innovazione',
     'Portare innovazione reale dove produce risultati misurabili. Prima il problema, poi il processo, poi il modello.'),
]

y = 4.5
for name, desc in pillars:
    tb(s, 3.8, y, 7.6, 0.9, name, 18, PL, bold=True)
    tb(s, 12.7, y, 19.1, 1.1, desc, 15, SL)
    y += 2.4

# ================================================================
# SLIDE 6 — 15 pilastri in 4 macroaree
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Il metodo tecnico \u2014 15 pilastri, 4 macroaree')
purple_bar(s)
tb(s, 3.0, 4.3, 27.9, 0.8, '15 pilastri tecnici organizzati in 4 macro-aree:', 16, SL)

rows = [
    ('COSTRUIRE',  'Architettura, Design, Qualit\u00e0',
     'Costruire solo la complessit\u00e0 necessaria', PD),
    ('VERIFICARE', 'Testing, Delivery, Osservabilit\u00e0',
     'Un sistema non \u00e8 governabile se non \u00e8 osservabile', PM),
    ('PROTEGGERE', 'Sicurezza, Dati, Scalabilit\u00e0',
     'Prima misuriamo. Poi ottimizziamo. Infine scaliamo', PL),
    ('GUIDARE',    'Team, Debito, AI, Decisioni, Costi, Comunicazione',
     'La tecnologia produce valore quando genera decisioni migliori', GREEN),
]

y = 5.8
for macro, pillars_str, principle, accent in rows:
    rect(s, 3.0, y, 27.4, 2.2, SD)
    rect(s, 3.0, y, 0.25, 2.2, accent)
    tb(s, 3.8, y+0.2, 7.1, 0.9, macro, 16, PL, bold=True)
    tb(s, 11.4, y+0.2, 11.4, 0.9, pillars_str, 14, SL)
    tb(s, 23.4, y+0.2, 7.5, 0.9, principle, 13, SM)
    y += 2.65

# ================================================================
# SLIDE 7 — Come lavorerei: 3 fasi
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Come lavorerei \u2014 tre fasi')
purple_bar(s)

phases = [
    ('0\u201360 gg', 'COMPRENDERE\nE INTERVENIRE', PD, [
        'Incontri con stakeholder',
        'Mappa sistemi, integrazioni, dati',
        'Mappa delle responsabilit\u00e0',
        'Identificazione rischi',
        'Primi interventi ad alto impatto',
        'Definizione metriche iniziali',
    ]),
    ('60\u2013180 gg', 'STABILIZZARE', PM, [
        'Ownership flussi critici',
        'Osservabilit\u00e0',
        'Gestione incidenti',
        'Processo decisionale',
        'Governance fornitori',
        'Roadmap condivisa',
        'Riduzione punti di dipendenza',
    ]),
    ('180+ gg', 'RENDERE\nAUTONOMI', PL, [
        'Trasferimento competenze',
        'Mentoring',
        'Documentazione essenziale',
        'Delega progressiva',
        'Rituali di governance',
        'Presidio strategico',
    ]),
]

cw = 9.9; gap = 0.65; x0 = 1.5; yt = 4.0; ch = 13.0
for i, (period, title, color, items) in enumerate(phases):
    col_card(s, x0 + i*(cw+gap), yt, cw, ch, period, title, color, items)

footer(s,
    "La presenza iniziale \u00e8 maggiore. Si riduce in funzione dell'autonomia raggiunta \u2014 non del tempo trascorso.",
    PL, 13)

# ================================================================
# SLIDE 8 — Percorso AI
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Il percorso AI \u2014 5 fasi')
purple_bar(s)

phases_ai = [
    ('1. Readiness',
     'Infrastruttura, qualit\u00e0 dati, maturit\u00e0 processi, cultura e competenze, governance e sicurezza.'),
    ("2. Casi d'uso",
     'Valore economico, frequenza problema, disponibilit\u00e0 dati, complessit\u00e0 integrazione, rischio, time-to-result.'),
    ('3. Esperimento misurabile',
     'Baseline \u2192 risultato atteso \u2192 utenti reali \u2192 KPI \u2192 criteri di valutazione \u2192 decisione finale: procedere, modificare o interrompere.'),
    ('4. Industrializzazione',
     'Integrazioni reali, controllo accessi, osservabilit\u00e0, gestione errori, versionamento, costi, fallback, eval periodiche.'),
    ('5. Adozione',
     "L'AI produce valore solo quando entra nei processi delle persone: formazione, linee guida, ownership, ROI monitoring."),
]

y = 4.3
for name, desc in phases_ai:
    rect(s, 3.0, y, 27.4, 2.0, SD)
    rect(s, 3.0, y, 0.25, 2.0, PL)
    tb(s, 3.8, y+0.2, 9.0, 0.8, name, 16, PL, bold=True)
    tb(s, 13.5, y+0.2, 17.0, 1.5, desc, 13, SL)
    y += 2.4

footer(s, "L'AI non parte dalla scelta dello strumento. Parte dalla maturit\u00e0 reale dell'organizzazione.", PL, 13)

# ================================================================
# SLIDE 9 — Principi AI
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'I miei principi AI')
purple_bar(s)

p_left = [
    ('Sistema, non demo',
     'Una demo pu\u00f2 impressionare. Un sistema deve essere affidabile, controllabile, integrato e mantenibile.'),
    ('Eval prima',
     "Prima di aumentare utenti o casi d'uso, definire come misurare la qualit\u00e0. Senza eval = slot machine."),
    ('Cost routing',
     'Non tutti i problemi richiedono il modello pi\u00f9 potente. Scegliere in funzione di compito, rischio e volume.'),
]
p_right = [
    ('Rischi espliciti',
     'Prompt injection, data leakage, allucinazioni, errori non rilevati, lock-in fornitore, costi incontrollati.'),
    ('Build il core, buy il resto',
     'Costruire il vantaggio competitivo. Acquistare la commodity. Ibrido per mantenere il valore differenziante.'),
    ('Responsabilit\u00e0 umana',
     "L'AI pu\u00f2 supportare una decisione. La responsabilit\u00e0 finale deve restare assegnata a una persona."),
]

y = 4.0
card_h = 3.2
gap_y = 0.5
for name, desc in p_left:
    rect(s, 2.5, y, 13.5, card_h, SD)
    rect(s, 2.5, y, 0.25, card_h, PD)
    tb(s, 3.2, y+0.2, 12.5, 0.8, name, 15, PL, bold=True)
    tb(s, 3.2, y+1.2, 12.5, 1.8, desc, 13, SL)
    y += card_h + gap_y

y = 4.0
for name, desc in p_right:
    rect(s, 17.5, y, 13.5, card_h, SD)
    rect(s, 17.5, y, 0.25, card_h, PM)
    tb(s, 18.2, y+0.2, 12.5, 0.8, name, 15, PL, bold=True)
    tb(s, 18.2, y+1.2, 12.5, 1.8, desc, 13, SL)
    y += card_h + gap_y

# ================================================================
# SLIDE 10 — Cosa cambia
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Cosa cambia')
purple_bar(s)

cols = [
    ('Per il CEO', PD, [
        'Maggiore controllo sugli investimenti',
        'Rischi visibili prima delle emergenze',
        'Roadmap collegata alla strategia',
        'Alternative comprensibili',
        'Responsabilit\u00e0 definite',
        'Migliore gestione fornitori',
        'Minore dipendenza da singole persone',
        'Maggiore prevedibilit\u00e0',
    ]),
    ('Per il responsabile tecnico', PM, [
        'Supporto nelle decisioni complesse',
        'Condivisione della pressione',
        'Priorit\u00e0 pi\u00f9 stabili',
        'Architettura pi\u00f9 chiara',
        'Maggiore forza nel dialogo col business',
        'Meno lavoro invisibile',
        'Crescita della leadership',
        'Maggiore autonomia del team',
    ]),
    ('Per il marketing', PL, [
        'Maggiore prevedibilit\u00e0',
        'Dati pi\u00f9 affidabili',
        'Meno dipendenza dalle urgenze',
        'Valutazione tecnica anticipata',
        'Migliore integrazione CRM/loyalty/e-comm',
        'Tempi e vincoli comunicati prima',
        'Automazioni con impatto misurabile',
        'Sperimentazione AI governata',
    ]),
]

cw = 9.9; gap = 0.65; x0 = 1.5; yt = 4.0; ch = 13.5
for i, (title, color, items) in enumerate(cols):
    col_card(s, x0 + i*(cw+gap), yt, cw, ch, '', title, color, items, item_size=13)

# ================================================================
# SLIDE 11 — Come misurare il risultato
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Come misurare il risultato')
purple_bar(s)

tb(s, 3.0, 4.3, 27.9, 0.8,
   'Le metriche definitive verranno concordate dopo una baseline iniziale.', 15, SL)

areas_l = [
    'Affidabilit\u00e0 dei flussi critici',
    'Numero e durata degli incidenti',
    'Tempo di rilevazione',
    'Tempo di ripristino',
    'Velocit\u00e0 dei rilasci',
    'Percentuale rilasci problematici',
    'Tempo per abilitare una campagna',
]
areas_r = [
    'Errori nelle integrazioni',
    'Qualit\u00e0 dei dati',
    'Attivit\u00e0 manuali eliminate',
    'Decisioni tecniche bloccate',
    'Dipendenza da singole persone',
    'Costi evitati o ottimizzati',
    'Valore prodotto dalle iniziative AI',
]

lines_l = [('OPERATIVIT\u00c0', 12, PL, True), ('', 4, SL, False)]
for a in areas_l: lines_l += [('\u2192  ' + a, 14, SL, False)]

lines_r = [('QUALIT\u00c0 & COSTO', 12, GREEN, True), ('', 4, SL, False)]
for a in areas_r: lines_r += [('\u2192  ' + a, 14, SL, False)]

multiline_tb(s, 3.0, 5.7, 13.5, 11.5, lines_l)
multiline_tb(s, 18.0, 5.7, 13.5, 11.5, lines_r)

footer(s,
    'Non prometto numeri prima di aver misurato il punto di partenza. Prometto chiarezza su cosa misurare e quale decisione deve supportare.',
    PL, 13)

# ================================================================
# SLIDE 12 — Security & Resilience Assessment — findings
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Security & Resilience Assessment — 23 luglio 2026')
purple_bar(s)

FIXED   = RGBColor(0x34, 0xD3, 0x99)
OPEN    = RGBColor(0xF8, 0x71, 0x71)
PARTIAL = RGBColor(0xFB, 0xBF, 0x24)

tb(s, 3.0, 4.3, 27.9, 0.75,
   'Assessment esterno, non autenticato, non invasivo — veralab.it e sottodomini pubblici.',
   13, SM, italic=True)

# findings: (id, title, severity, severity_color, description)
findings_sec = [
    ('SEC-01', 'Login B2B via HTTP',
     'ALTA', OPEN,
     'b2b.veralab.it raggiungibile e funzionante in HTTP. Credenziali e cookie in chiaro. HSTS assente sul sottodominio.'),
    ('SEC-02', 'Cookie sessione B2B privi di Secure/HttpOnly/SameSite',
     'ALTA', OPEN,
     'PHPSESSID e last_session_id senza attributi di sicurezza. Cookie inviabili in HTTP. IP e timestamp memorizzati lato client.'),
    ('SEC-03', 'Anti-clickjacking assente sul frontend headless',
     'MEDIA', PARTIAL,
     'Homepage, login e Magazine privi di frame-ancestors e X-Frame-Options. Checkout Shopify è invece correttamente protetto.'),
    ('SEC-04', 'Content Security Policy assente',
     'MEDIA', PARTIAL,
     'Nonce già generati server-side ma CSP non inviata. Nessuna difesa aggiuntiva contro XSS e script di terze parti.'),
    ('SEC-05', 'CORS permissivo su api2.veralabtech.net',
     'MEDIA\nda validare', PARTIAL,
     'Origin arbitraria riflessa con Allow-Credentials: true. Sfruttabile se endpoint espongono dati autenticati.'),
]
findings_rel = [
    ('REL-01', 'Route inesistenti restituiscono 500',
     'MEDIA', PARTIAL,
     '/.well-known/security.txt e route /it-it/* casuali → 500. Collection inesistenti → soft 404 con redirect a homepage.'),
    ('REL-02', '/blogs/news — rotta legacy in errore',
     'MEDIA', PARTIAL,
     'Non linkata dalla navigazione corrente. Redirect mancante verso /it-it/magazine/. Possibile impatto su backlink storici.'),
    ('REL-03', 'Payload HTML molto grandi e pagine non cacheabili',
     'MEDIA/ALTA\nbusiness', PARTIAL,
     'Homepage ~775 KB, Magazine ~1,67 MB. oxygen-full-page-cache: uncacheable. TTFB osservato 4-8s. Rischio su picchi e lanci.'),
    ('CFG-01', 'Cookie con durata 82 anni',
     'BASSA', YELLOW,
     'Max-Age=2592000000 (probabilmente ms invece di sec). Interessa country, language, cartId. Attributi Secure/HttpOnly/SameSite OK.'),
]

# Left column — Security
tb(s, 2.5, 5.3, 15.0, 0.65, 'SICUREZZA', 11, OPEN, bold=True)
rh = 2.2; gap = 0.3; y = 6.1
for fid, title, sev, sev_clr, desc in findings_sec:
    rect(s, 2.5, y, 15.0, rh, SD)
    rect(s, 2.5, y, 0.18, rh, sev_clr)
    tb(s, 2.85, y+0.12, 2.2, 0.6, fid, 10, sev_clr, bold=True)
    tb(s, 5.2,  y+0.12, 12.0, 0.75, title, 12, WHITE, bold=True)
    tb(s, 2.85, y+0.9,  14.5, 1.1, desc, 11, SL)
    y += rh + gap

# Right column — Resilienza + Config
tb(s, 18.3, 5.3, 14.5, 0.65, 'RESILIENZA & CONFIGURAZIONE', 11, PARTIAL, bold=True)
y = 6.1
for fid, title, sev, sev_clr, desc in findings_rel:
    rect(s, 18.3, y, 14.5, rh, SD)
    rect(s, 18.3, y, 0.18, rh, sev_clr)
    tb(s, 18.65, y+0.12, 2.2, 0.6, fid, 10, sev_clr, bold=True)
    tb(s, 21.0,  y+0.12, 11.5, 0.75, title, 12, WHITE, bold=True)
    tb(s, 18.65, y+0.9,  13.8, 1.1, desc, 11, SL)
    y += rh + gap

footer(s,
    'Nessuna evidenza di compromissione rilevata. Priorità assoluta: portale B2B — credenziali mai in HTTP.',
    OPEN, 13)

# ================================================================
# SLIDE 13 — Cosa funziona / cosa è già corretto
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Cosa funziona già — controlli con esito positivo')
purple_bar(s)

tb(s, 3.0, 4.3, 27.9, 0.75,
   "L'obiettività richiede di citare anche ciò che è configurato correttamente.",
   13, SM, italic=True)

positives = [
    ('TLS / HTTPS', [
        'SSL Labs: A+',
        'Solo TLS 1.2 e 1.3',
        'Redirect HTTP → HTTPS funzionante',
        'www → apex redirect corretto',
        'HSTS max-age=31536000',
    ]),
    ('Cookie principale', [
        'HttpOnly; Secure; SameSite=Lax',
        'X-Content-Type-Options: nosniff',
        'TRACE disabilitato (405)',
        'Nessuno stack trace visibile nei 500',
        'CDN Cloudflare attivo',
    ]),
    ('Checkout & landing Shopify', [
        'Redirect HTTP → HTTPS',
        'X-Frame-Options: DENY',
        'CSP con frame-ancestors none',
        'noindex su pagine sensibili',
        'X-Content-Type-Options: nosniff',
    ]),
    ('Frontend & asset', [
        'Source map .js → 404 (non esposti)',
        'Parametri di ricerca codificati correttamente',
        'Dominio tecnico Shopify: richiede autenticazione',
        'Magazine operativo su /it-it/magazine/',
        'JSON-LD OnlineStore presente',
    ]),
]

cw = 7.8; gap = 0.55; x0 = 1.5; yt = 5.5; ch = 11.5
for i, (title, items) in enumerate(positives):
    col_card(s, x0 + i*(cw+gap), yt, cw, ch, '', title, GREEN, items, item_size=12)

footer(s,
    'La base è solida. I problemi rilevati sono reali ma risolvibili — non indicano la necessità di sostituire Shopify.',
    GREEN, 13)

# ================================================================
# SLIDE 14 — Priorità di intervento
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Priorità di intervento')
purple_bar(s)

p24h = [
    'Forzare HTTPS su tutto il portale B2B',
    'Correggere attributi cookie B2B (Secure, HttpOnly, SameSite)',
    'Attivare HSTS sul sottodominio B2B',
    'Verificare link/email/QR che puntano al B2B in HTTP',
    'Correggere redirect /blogs/news → /it-it/magazine/',
    'Validare configurazione CORS su api2.veralabtech.net',
]
p7d = [
    'Route inesistenti → sempre 404, mai 500',
    'Eliminare soft 404 e redirect generici a homepage',
    'Introdurre frame-ancestors + X-Frame-Options',
    'Attivare CSP in modalità Report-Only',
    'Configurare alerting su error ratio 5xx per route',
    'Correggere durata cookie (2.592.000 sec, non ms)',
    'Pubblicare /.well-known/security.txt',
]
p30d = [
    'CSP in enforcement',
    'Ridurre payload HTML e dati serializzati',
    'Introdurre caching edge (stale-while-revalidate)',
    'Real User Monitoring + Core Web Vitals',
    'Performance budget + Security Definition of Done',
    'Penetration test autenticato: account, B2B, checkout, API',
]

cw = 9.9; gap = 0.65; x0 = 1.5; yt = 4.0; ch = 13.5
col_card(s, x0,          yt, cw, ch, 'ENTRO 24 ORE', 'B2B — ALTA PRIORITÀ', OPEN,   p24h, 13)
col_card(s, x0+cw+gap,   yt, cw, ch, 'ENTRO 7 GIORNI', 'HARDENING & ROUTING', PARTIAL, p7d, 13)
col_card(s, x0+2*(cw+gap), yt, cw, ch, 'ENTRO 30 GIORNI', 'GOVERNANCE & PENTEST', PL, p30d, 13)

footer(s,
    'La sequenza riflette impatto e reversibilità — non complessità tecnica. Nessuno di questi interventi richiede un refactoring del frontend.',
    SM, 13)

# ================================================================
# SLIDE 15 — Closing
# ================================================================
s = prs.slides.add_slide(BL); bg(s)
title_bar(s, 'Il prossimo passo')
purple_bar(s)

# Big quote box
rect(s, 2.5, 4.3, 27.9, 4.2, SD)
rect(s, 2.5, 4.3, 0.3, 4.2, PD)
tb(s, 3.3, 4.6, 26.5, 3.6,
   '"Entro con un coinvolgimento operativo forte, porto valore immediato sulle priorit\u00e0 che contano, costruisco il sistema di governance e accompagno team e azienda verso una progressiva autonomia."',
   17, PL, italic=True)

tb(s, 3.0, 9.3, 8.0, 0.7, 'PRIMO PASSO', 12, SM, bold=True)
lines = [
    ('Una sessione con CEO, responsabile tecnico e marketing per definire:', 15, SL, False),
    ('', 5, SL, False),
    ('1.  i risultati aziendali prioritari', 15, SL, False),
    ('2.  i flussi tecnologici pi\u00f9 critici', 15, SL, False),
    ('3.  le decisioni oggi bloccate', 15, SL, False),
    ('4.  il primo risultato visibile', 15, SL, False),
    ('5.  responsabilit\u00e0 e modalit\u00e0 di collaborazione', 15, SL, False),
]
multiline_tb(s, 3.0, 10.2, 27.9, 7.0, lines)
footer(s, '108 Vision \u2014 Costruiamo la direzione, non solo il codice.', SM, 13)

# ── Save ─────────────────────────────────────────────────────
out = r'c:\Code\Documents\Lavoro\Personale\Vision\tracks\108-cto\PRES_Veralab_FCTO_Operativo_v2.pptx'
prs.save(out)
print(f'OK: {out}')
