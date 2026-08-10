"""
Script: create_pptx_veralab_fcto.py
Output: PRES_Veralab_FCTO_EliosScoglio.pptx

Struttura isomorfa a PRES_EticaSoluzioni_EliosScoglio.pptx:
  Slide 1  — Cover
  Slide 2  — Chi sono
  Slide 3  — Cosa ho trovato (SPECIFICO VERALAB — sostituisce slide "Etica" dedicata)
  Slide 4  — I gap architetturali (inferiti — marcati come probabili)
  Slide 5  — Cosa faccio concretamente
  Slide 6  — Risultati dimostrabili
  Slide 7  — Il mio approccio: due filoni
  Slide 8  — Il metodo: Know-How → Step → KPI
  Slide 9  — Modello Fractional CTO
  Slide 10 — Come funziona in pratica
  Slide 11 — Perche sono rilevante per VOI
  Slide 12 — Prossimo passo: Tech Assessment
  Slide 13 — Chiusura
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    from pptx.enum.shapes import MSO_SHAPE_TYPE as MSO_SHAPE

# ─── PALETTE 108 Vision ───────────────────────────────────────────────────────
INK_950    = RGBColor(0x0F, 0x17, 0x2A)
INK_900    = RGBColor(0x1E, 0x29, 0x3B)
INK_800    = RGBColor(0x33, 0x41, 0x55)
INK_200    = RGBColor(0xE2, 0xE8, 0xF0)
INK_50     = RGBColor(0xF8, 0xFA, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
VIOLET_700 = RGBColor(0x6D, 0x28, 0xD9)
VIOLET_500 = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_400 = RGBColor(0xA7, 0x8B, 0xFA)
GREEN_400  = RGBColor(0x34, 0xD3, 0x99)
RED_400    = RGBColor(0xF8, 0x71, 0x71)
AMBER_400  = RGBColor(0xFB, 0xBF, 0x24)

# ─── PRESENTAZIONE ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = INK_950
    return slide


def text(slide, left, top, width, height, txt, size=18, bold=False,
         color=None, align=PP_ALIGN.LEFT, italic=False):
    if color is None:
        color = WHITE
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = 'Inter'
    return tf


def add_para(tf, txt, size=14, bold=False, color=None, space_before=6, italic=False):
    if color is None:
        color = INK_200
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = 'Inter'
    p.space_before = Pt(space_before)
    return p


def bar(slide, top, height=0.8):
    """Violet accent bar on the left."""
    s = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(top), Inches(0.08), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = VIOLET_700
    s.line.fill.background()


def hline(slide, left, top, width, color=None, thickness=0.04):
    if color is None:
        color = VIOLET_700
    s = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(thickness))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def card(slide, x, y, w, h, border_color=None):
    if border_color is None:
        border_color = VIOLET_700
    s = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = INK_900
    s.line.color.rgb = border_color
    s.line.width = Pt(1.5)
    return s


def severity_dot(slide, x, y, color):
    s = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        Inches(x), Inches(y), Inches(0.18), Inches(0.18))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()

text(slide, 1.0, 1.2, 11, 1.2, 'ELIOS SCOGLIO', size=48, bold=True)
text(slide, 1.0, 2.45, 11, 0.7,
     'Software & Architecture Manager  |  Fractional CTO',
     size=22, color=VIOLET_400)
hline(slide, 1.0, 3.3, 3.5)
text(slide, 1.0, 3.6, 11, 0.6,
     'Proposta Fractional CTO — Veralab / Re-Forme SRL',
     size=18, color=INK_200)
text(slide, 1.0, 4.25, 11, 0.5,
     'Ho analizzato il vostro stack prima di arrivare qui.',
     size=16, color=INK_200, italic=True)
text(slide, 1.0, 6.2, 11, 0.5,
     '108 Vision — Costruiamo la direzione, non solo il codice.',
     size=13, color=INK_800)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CHI SONO
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Chi sono', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.75, 11, 0.45,
     'Software & Architecture Manager — piattaforme mission-critical da 10+ anni.',
     size=16, color=INK_200)

exp = [
    ('30M transazioni/anno', 'Zero-downtime, spike gestiti, inventory real-time su sistemi distribuiti'),
    ('93 componenti, 7 livelli', '3 team di sviluppo coordinati sulla stessa roadmap tecnica'),
    ('Compliance stringente', 'GDPR quotidiano, normativa fiscale, integrazioni enti statali'),
    ('Legacy modernizzato', 'Da CORBA 23 anni a microservizi gRPC — senza fermare la produzione'),
]
for i, (title, desc) in enumerate(exp):
    y = 2.35 + i * 0.88
    text(slide, 1.5, y, 4.0, 0.4, title, size=16, bold=True, color=VIOLET_400)
    text(slide, 5.7, y, 7.2, 0.5, desc,  size=15, color=INK_200)

hline(slide, 1.2, 5.9, 11.0)

results = [
    ('Deploy freq.',  '+400%'),
    ('Durata deploy', '-91%'),
    ('Costo sviluppo (AI)', '-77%/-82%'),
    ('Team satisfaction', '+50%'),
]
for i, (label, number) in enumerate(results):
    x = 1.4 + i * 3.0
    text(slide, x, 6.05, 2.8, 0.35, label,  size=12, color=INK_800)
    text(slide, x, 6.45, 2.8, 0.5,  number, size=22, bold=True, color=VIOLET_400)

text(slide, 1.2, 7.05, 11, 0.3,
     'Non scrivo codice. Faccio in modo che le decisioni tecniche siano quelle giuste.',
     size=13, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — COSA HO TROVATO (specifico Veralab — dati live)
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75,
     'Cosa ho trovato — prima di questa call', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.72, 11, 0.4,
     'Analisi live su veralab.it con script su endpoint reali — non tool generici.',
     size=15, color=INK_200, italic=True)

# Tabella finding
headers = ['Finding', 'Severity', 'Stato']
col_x = [1.2, 7.8, 10.2]
col_w = [6.4, 2.2, 2.8]
for j, h in enumerate(headers):
    text(slide, col_x[j], 2.25, col_w[j], 0.32,
         h, size=12, bold=True, color=INK_800)
hline(slide, 1.2, 2.6, 11.0, INK_800)

findings = [
    ('/blogs/news  HTTP 500', 'CRITICO', 'In produzione adesso', RED_400),
    ('Nessun alerting su errori 5xx', 'ALTO', 'Non sapete quando si rompe', RED_400),
    ('Content-Security-Policy assente', 'ALTO', 'XSS via app di terze parti', RED_400),
    ('Structured data JSON-LD assente', 'MEDIO', 'Nessun rich result su Google', AMBER_400),
    ('Homepage 765 KB — LCP stimato 3-4s', 'MEDIO', 'Core Web Vitals sotto soglia', AMBER_400),
    ('X-Frame-Options assente', 'MEDIO', 'Clickjacking su checkout', AMBER_400),
    ('Overskin — placeholder Aruba, non live', 'INFO', 'Secondo brand senza infrastruttura', INK_800),
]
for i, (finding, sev, stato, col) in enumerate(findings):
    y = 2.72 + i * 0.54
    severity_dot(slide, col_x[0] - 0.3, y + 0.08, col)
    text(slide, col_x[0], y, col_w[0], 0.45, finding, size=13, color=WHITE)
    text(slide, col_x[1], y, col_w[1], 0.45, sev,     size=12, bold=True, color=col)
    text(slide, col_x[2], y, col_w[2], 0.45, stato,   size=12, color=INK_200)

text(slide, 1.0, 6.65, 11, 0.45,
     'Il punto piu urgente: /blogs/news risponde HTTP 500 — il Magazine e rotto, '
     'non e nei vostri alert, probabilmente da settimane.',
     size=12, bold=True, color=RED_400)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — I GAP ARCHITETTURALI (inferiti — marcati come probabili)
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75,
     'I gap architetturali (da verificare in onboarding)', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.72, 11, 0.45,
     'Inferiti dalla struttura del business — non misurabili dall\'esterno. '
     'Li presento come ipotesi, non come certezze.',
     size=14, color=AMBER_400, italic=True)

gaps = [
    (
        'Inventory non unificato',
        '[probabile]',
        'Shopify (stock online) e NAV (14 store fisici + B2B) quasi certamente non condividono '
        'un single source of truth real-time. Conseguenze: click & collect bloccato, '
        'overselling sui lanci, loyalty cross-channel limitata.',
        AMBER_400,
    ),
    (
        'Decisione ERP aperta',
        '[probabile]',
        'Dynamics NAV fuori mainstream support su versioni pre-2018. '
        'Migrazione a Business Central: 6-18 mesi, 100-400K EUR. '
        'Ogni trimestre senza una decisione informata aumenta il rischio.',
        AMBER_400,
    ),
    (
        'Architettura multi-brand non definita',
        '[probabile]',
        'Veralab + Overskin su Shopify standard: senza Shopify Plus o scelta architetturale '
        'esplicita, nascera debito tecnico impossibile da sciogliere.',
        AMBER_400,
    ),
    (
        'Monitoring assente',
        '[verificato]',
        'Il 500 del Magazine non e nei vostri alert. '
        'Non avete visibilita su cosa succede in produzione.',
        RED_400,
    ),
]
for i, (title, tag, desc, col) in enumerate(gaps):
    y = 2.4 + i * 1.12
    card(slide, 1.2, y, 10.9, 1.0, col)
    text(slide, 1.55, y + 0.08, 4.5, 0.38,
         title, size=15, bold=True, color=col)
    text(slide, 6.2, y + 0.08, 1.4, 0.38,
         tag, size=11, italic=True, color=INK_800)
    text(slide, 1.55, y + 0.5, 10.2, 0.4,
         desc, size=12, color=INK_200)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — COSA FACCIO CONCRETAMENTE
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Cosa faccio concretamente', size=34, bold=True)
bar(slide, 1.3)

rows = [
    ('Strategia tecnica',  'Roadmap, decisioni architetturali, allineamento business ↔ tech'),
    ('Architettura',       'Decisioni strutturali con trade-off espliciti: NAV vs BC, Shopify Plus vs headless, '
                           'inventory sync approach, architettura Overskin'),
    ('Team',               'Sviluppo del Tech Lead interno. Costruisco autonomia — non dipendenza da me'),
    ('Stakeholder',        'Traduco il tecnico in business: impatto in EUR, rischi concreti, timeline realistiche'),
    ('Monitoraggio',       'Visibilita su produzione: alert, metriche, 4 Golden Signals'),
]
hline(slide, 1.2, 1.7, 11.0, INK_800)
text(slide, 1.3, 1.75, 4.0, 0.3, 'Responsabilita', size=12, bold=True, color=INK_800)
text(slide, 5.5, 1.75, 7.5, 0.3, 'Cosa significa in pratica', size=12, bold=True, color=INK_800)
hline(slide, 1.2, 2.1, 11.0, INK_800)

for i, (resp, desc) in enumerate(rows):
    y = 2.2 + i * 0.88
    text(slide, 1.3, y, 4.0, 0.75, resp, size=15, bold=True, color=VIOLET_400)
    text(slide, 5.5, y, 7.5, 0.75, desc, size=14, color=INK_200)
    if i < len(rows) - 1:
        hline(slide, 1.2, y + 0.82, 11.0, INK_800)

text(slide, 1.2, 6.65, 11, 0.45,
     'Cosa NON faccio: scrivere codice, gestire sprint, fare il PM, risolvere bug in autonomia.',
     size=13, bold=True, color=RED_400)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RISULTATI DIMOSTRABILI
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Risultati dimostrabili', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.72, 11, 0.4,
     'Questi numeri vengono da sistemi reali in produzione — non da POC o demo.',
     size=15, color=INK_200, italic=True)

metrics = [
    ('Deploy frequency',        'Da ogni 6 settimane',  'Ogni settimana',  '+400%', VIOLET_400),
    ('Durata deploy',           'Da 4 ore',             '22 minuti',       '-91%',  GREEN_400),
    ('Costo sviluppo (AI)',     'Baseline',             'Task specifici',  '-77%/-82%', GREEN_400),
    ('Tempo analisi requisiti', 'Da 2-3 giorni',        '2 ore',           '-92%',  GREEN_400),
    ('Team satisfaction',       'Baseline',             '',                '+50%',  VIOLET_400),
    ('Sprint velocity',         'Baseline',             '',                '+30%',  VIOLET_400),
]

hline(slide, 1.2, 2.22, 11.0, INK_800)
for j, h in enumerate(['Metrica', 'Prima', 'Dopo', 'Miglioramento']):
    xs = [1.2, 4.5, 7.2, 10.2]
    ws = [3.0, 2.5, 2.8, 2.5]
    text(slide, xs[j], 2.28, ws[j], 0.3, h, size=12, bold=True, color=INK_800)
hline(slide, 1.2, 2.62, 11.0, INK_800)

for i, (metric, before, after, delta, col) in enumerate(metrics):
    y = 2.72 + i * 0.68
    text(slide, 1.2,  y, 3.0, 0.5, metric, size=14, bold=True, color=WHITE)
    text(slide, 4.5,  y, 2.5, 0.5, before, size=13, color=INK_800)
    text(slide, 7.2,  y, 2.8, 0.5, after,  size=13, color=INK_200)
    text(slide, 10.2, y, 2.5, 0.5, delta,  size=16, bold=True, color=col)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — IL MIO APPROCCIO: DUE FILONI
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Il mio approccio: due filoni paralleli', size=34, bold=True)
bar(slide, 1.3)

# Filone 1 — sinistra
card(slide, 1.0, 1.65, 5.5, 5.45, VIOLET_700)
text(slide, 1.3, 1.82, 5.0, 0.4, 'FILONE 1 — PRESENTE', size=13, bold=True, color=VIOLET_700)
text(slide, 1.3, 2.22, 5.0, 0.55, 'Concretizzare e ottimizzare', size=18, bold=True)
tf1 = text(slide, 1.3, 2.85, 5.0, 4.0, '', size=14, color=INK_200)
f1_items = [
    'Fix Magazine + monitoring su 5xx (primo giorno)',
    'Security headers via Cloudflare (CSP, X-Frame)',
    'Structured data JSON-LD → rich results Google',
    'Performance homepage (765 KB → target < 200 KB)',
    'Pipeline di deploy stabile, code review strutturata',
    'Incident playbook: cosa fare quando il checkout\nva giu durante un lancio influencer',
]
for item in f1_items:
    add_para(tf1, f'→  {item}', size=13, color=INK_200, space_before=8)
text(slide, 1.3, 6.65, 5.0, 0.45,
     'Zero rivoluzioni. Valore visibile in 2-4 settimane.',
     size=12, bold=True, color=VIOLET_400)

# Filone 2 — destra
card(slide, 6.85, 1.65, 5.5, 5.45, VIOLET_500)
text(slide, 7.15, 1.82, 5.0, 0.4, 'FILONE 2 — VISIONE', size=13, bold=True, color=VIOLET_500)
text(slide, 7.15, 2.22, 5.0, 0.55, 'Dove andare nei prossimi 12-24 mesi', size=18, bold=True)
tf2 = text(slide, 7.15, 2.85, 5.0, 4.0, '', size=14, color=INK_200)
f2_items = [
    'Inventory unificato: Shopify ↔ NAV',
    'NAV → Business Central: go/no-go con dati',
    'Architettura Overskin (prima del debito)',
    'Loyalty VERABILIA davvero cross-channel',
    'Profilo cliente unificato: CRM/CDP',
    'AI dove ha senso: predizione stock,\npersonalizzazione, chatbot post-vendita',
]
for item in f2_items:
    add_para(tf2, f'→  {item}', size=13, color=INK_200, space_before=8)
text(slide, 7.15, 6.65, 5.0, 0.45,
     'Ogni decisione ha costo, beneficio e piano di rollback.',
     size=12, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — IL METODO: KNOW-HOW → STEP → KPI
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Il metodo: Know-How → Step → KPI', size=34, bold=True)
bar(slide, 1.3)

phases = [
    ('MESE 1', 'ASCOLTO', VIOLET_700, [
        'Non cambio niente in produzione.',
        'Mappo: Shopify, NAV, CRM, WMS, POS, loyalty.',
        'Misuro la baseline di tutti i KPI.',
        'Output: "State of the Stack" + Piano 90gg',
        'Fix immediato: /blogs/news + alert 5xx',
    ]),
    ('MESI 2-3', 'DECISIONI CRITICHE', VIOLET_500, [
        'Assessment NAV → BC: go/no-go con dati.',
        'Design inventory sync Shopify ↔ NAV.',
        'Architettura Overskin (se in corso).',
        'ADR per le prime 5 decisioni arch.',
    ]),
    ('MESI 4-6', 'OMNICHANNEL', VIOLET_400, [
        'Inventory unificato (source of truth).',
        'Click & collect su tutti i 14 store.',
        'Loyalty VERABILIA cross-channel.',
        'Profilo cliente online + offline.',
    ]),
]

for i, (period, title, col, items) in enumerate(phases):
    x = 0.7 + i * 4.2
    card(slide, x, 1.7, 3.9, 4.85, col)
    text(slide, x + 0.2, 1.85, 3.5, 0.35, period, size=11, bold=True, color=col)
    text(slide, x + 0.2, 2.25, 3.5, 0.55, title,  size=18, bold=True)
    tf = text(slide, x + 0.2, 2.88, 3.5, 3.3, '', size=13, color=INK_200)
    for item in items:
        add_para(tf, item, size=13, color=INK_200, space_before=9)

text(slide, 1.0, 6.85, 11, 0.35,
     'Non prometto numeri prima di misurare il baseline. Prima misuro, poi prometto.',
     size=13, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MODELLO FRACTIONAL CTO
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Perche il modello Fractional CTO', size=34, bold=True)
bar(slide, 1.3)

headers2 = ['', 'Full-time (Tech Lead / Dir. Tecnico)', 'Fractional CTO']
col_x2 = [1.2, 3.5, 8.5]
col_w2 = [2.1, 4.8, 4.6]
for j, h in enumerate(headers2):
    text(slide, col_x2[j], 1.7, col_w2[j], 0.35,
         h, size=13, bold=True,
         color=INK_800 if j == 1 else VIOLET_400 if j == 2 else INK_200)
hline(slide, 1.2, 2.1, 11.0, INK_800)

rows2 = [
    ('Costo annuo',         '120-180K EUR all-in',        '84-96K EUR (3gg/sett standard)'),
    ('Rischio hiring',      '6 mesi per capire se funziona', 'Trial via Tech Assessment — output prima'),
    ('Prospettiva',         'Solo interna',                'Esterna + interna — cross-pollination'),
    ('Visione settore',     'Un settore = visione limitata', 'Multi-settore: ticketing, e-commerce, retail, PA'),
    ('Exit',                'Se esce = panico + 6 mesi vuoto', 'Exit pianificata → team autonomo come deliverable'),
]
for i, (label, full, frac) in enumerate(rows2):
    y = 2.22 + i * 0.85
    text(slide, col_x2[0], y, col_w2[0], 0.7, label, size=14, bold=True, color=VIOLET_400)
    text(slide, col_x2[1], y, col_w2[1], 0.7, full,  size=13, color=INK_800)
    text(slide, col_x2[2], y, col_w2[2], 0.7, frac,  size=13, color=GREEN_400)
    if i < len(rows2) - 1:
        hline(slide, 1.2, y + 0.79, 11.0, INK_900)

text(slide, 1.0, 6.55, 11, 0.55,
     'Il vantaggio chiave: stessa seniority enterprise a costo proporzionale, '
     'con la possibilita di scalare su o giu ogni trimestre.',
     size=14, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COME FUNZIONA IN PRATICA
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Come funziona in pratica', size=34, bold=True)
bar(slide, 1.3)

details = [
    ('Presenza standard',     '3 giorni/settimana (remoto) + 1-2 presenze mensili a Milano'),
    ('Bridge iniziale',       '4 giorni/sett per max 4 mesi — piu presenza nelle prime decisioni (NAV, Overskin)'),
    ('Sessioni mensili fisse','Strategic Planning 2h (CEO), Architecture Review 2h (Tech Lead), Team Mentoring 1h'),
    ('Giorni off',            'Disponibilita async risposta entro 4h — emergenze entro 2h'),
    ('Deliverable mensili',   'Report scritto, ADR, roadmap aggiornata, metriche DORA'),
    ('Commitment minimo',     'Tech Assessment → poi 6 mesi (review trimestrale)'),
    ('Entry point',           'Tech Assessment (3 giorni) — output concreto, zero vincolo'),
]

hline(slide, 1.2, 1.7, 11.0, INK_800)
for j, h in enumerate(['Aspetto', 'Dettaglio']):
    xs = [1.3, 4.8]
    text(slide, xs[j], 1.76, 3.0, 0.3, h, size=12, bold=True, color=INK_800)
hline(slide, 1.2, 2.1, 11.0, INK_800)

for i, (asp, det) in enumerate(details):
    y = 2.2 + i * 0.65
    text(slide, 1.3, y, 3.3, 0.55, asp, size=14, bold=True, color=VIOLET_400)
    text(slide, 4.8, y, 7.9, 0.55, det, size=13, color=INK_200)
    if i < len(details) - 1:
        hline(slide, 1.2, y + 0.6, 11.0, INK_900)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PERCHE SONO RILEVANTE PER VOI
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Perche sono rilevante per Veralab in particolare', size=32, bold=True)
bar(slide, 1.3)

text(slide, 1.5, 1.72, 5.0, 0.32, 'IL VOSTRO CONTESTO', size=12, bold=True, color=INK_800)
text(slide, 8.0, 1.72, 5.0, 0.32, 'LA MIA ESPERIENZA DIRETTA', size=12, bold=True, color=VIOLET_400)
hline(slide, 1.2, 2.08, 11.2, INK_800)

matches = [
    ('Lanci prodotto con spike traffico (influencer)',
     'Gestisco on-sale con spike prevedibili — un\'ora di downtime = decine di K EUR persi'),
    ('Inventory multi-canale (14 store + online)',
     'Sistemi con inventory real-time e accesso concorrente su stock critico'),
    ('Shopify + NAV — integrazione eterogenea',
     'Integrazioni esterne con circuit breaker, retry, idempotenza — conosco il pattern'),
    ('Loyalty cross-channel (VERABILIA)',
     'Profili cliente unificati su sistemi eterogenei — problema di architettura, non di feature'),
    ('Secondo brand da lanciare (Overskin)',
     'Architetture multi-brand e multi-tenant su piattaforme condivise'),
    ('Monitoring assente',
     'Observability come governance: alert, metriche, Golden Signals su ogni servizio'),
]
for i, (their, mine) in enumerate(matches):
    y = 2.22 + i * 0.72
    text(slide, 1.5, y, 6.1, 0.58, their, size=13, color=INK_200)
    text(slide, 8.0, y, 5.1, 0.58, mine,  size=13, bold=True, color=VIOLET_400)
    if i < len(matches) - 1:
        hline(slide, 1.2, y + 0.66, 11.2, INK_900)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PROSSIMO PASSO: TECH ASSESSMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Prossimo passo: Tech Assessment', size=34, bold=True)
bar(slide, 1.3)

# Box principale
card(slide, 1.2, 1.65, 10.7, 3.35, VIOLET_700)
text(slide, 1.6, 1.82, 10.0, 0.5,
     'Tech Assessment — 3 giorni, output concreto, zero vincolo',
     size=18, bold=True, color=VIOLET_400)

tf_out = text(slide, 1.6, 2.42, 5.0, 2.3, '', size=14, color=INK_200)
outputs = [
    'Mappa completa: Shopify, NAV, CRM, WMS, POS, loyalty',
    'Integration map: dove i dati si rompono tra sistemi',
    'Top 5 rischi con impatto stimato in EUR',
    'Decisione NAV → Business Central: go/no-go con dati',
    'Architettura Overskin: raccomandazione con trade-off',
]
for o in outputs:
    add_para(tf_out, f'→  {o}', size=13, color=GREEN_400, space_before=7)

tf_cond = text(slide, 7.0, 2.42, 4.6, 2.3, '', size=14, color=INK_200)
conds = [
    ('Roadmap tecnica 12 mesi prioritizzata', WHITE),
    ('Fix del Magazine (primo giorno)', WHITE),
    ('', WHITE),
    ('Costo: EUR 2.500 – 3.500 (una tantum)', VIOLET_400),
    ('Se proseguite: detratto dal primo mese', GREEN_400),
    ('Se non proseguite: avete la roadmap', GREEN_400),
]
for txt_c, col_c in conds:
    add_para(tf_cond, txt_c, size=13, color=col_c, space_before=7)

# Box poi
card(slide, 1.2, 5.15, 10.7, 1.1, VIOLET_500)
text(slide, 1.6, 5.32, 10.0, 0.4,
     'Poi: 6 mesi (commitment minimo) con obiettivi concordati e KPI misurabili.',
     size=15, color=INK_200)
text(slide, 1.6, 5.75, 10.0, 0.35,
     'Review a fine ogni trimestre: continuare, scalare, o chiudere.',
     size=14, color=INK_800)

text(slide, 1.0, 6.55, 11, 0.55,
     'Il rischio non e provare il Fractional CTO. '
     'Il rischio e continuare senza direzione tecnica per altri 12 mesi.',
     size=14, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CHIUSURA
# ═══════════════════════════════════════════════════════════════════════════════
slide = dark_slide()

hline(slide, 3.5, 2.8, 6.3)
text(slide, 1.0, 1.4, 11, 1.1,
     '"Non vendo ore.\nVendo direzione e risultati misurabili.',
     size=28, bold=True, align=PP_ALIGN.CENTER)
text(slide, 1.0, 2.95, 11, 0.7,
     'Il giorno dopo l\'onboarding, il Magazine e online.',
     size=20, color=VIOLET_400, align=PP_ALIGN.CENTER)
text(slide, 1.0, 3.65, 11, 0.7,
     'Il mese dopo il Tech Assessment, i numeri parlano — non le slide."',
     size=20, color=INK_200, align=PP_ALIGN.CENTER, italic=True)
hline(slide, 3.5, 4.6, 6.3)

text(slide, 1.0, 5.05, 11, 0.55,
     'Elios Scoglio  |  elios@108vision.it  |  108vision.it',
     size=15, color=INK_200, align=PP_ALIGN.CENTER)
text(slide, 1.0, 5.65, 11, 0.45,
     'Software & Architecture Manager  |  Fractional CTO — 108 Vision',
     size=13, color=INK_800, align=PP_ALIGN.CENTER)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
out_path = (r'c:\Code\Documents\Lavoro\Personale\Vision\tracks\108-cto'
            r'\PRES_Veralab_FCTO_EliosScoglio.pptx')
prs.save(out_path)
print(f'PowerPoint salvato: {out_path}')
