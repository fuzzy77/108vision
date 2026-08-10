from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 108 Vision palette
INK_950 = RGBColor(0x0F, 0x17, 0x2A)
INK_900 = RGBColor(0x1E, 0x29, 0x3B)
INK_800 = RGBColor(0x33, 0x41, 0x55)
INK_200 = RGBColor(0xE2, 0xE8, 0xF0)
INK_50  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
VIOLET_700 = RGBColor(0x6D, 0x28, 0xD9)
VIOLET_500 = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_400 = RGBColor(0xA7, 0x8B, 0xFA)
GREEN_400  = RGBColor(0x34, 0xD3, 0x99)
RED_400    = RGBColor(0xF8, 0x71, 0x71)
AMBER_400  = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = INK_950
    return slide

def text(slide, left, top, width, height, txt, size=18, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = 'Inter'
    return tf

def add_para(tf, txt, size=14, bold=False, color=None, space_before=6):
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color if color else INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(space_before)
    return p

def bar(slide, top):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(0.08), Inches(0.8))
    s.fill.solid()
    s.fill.fore_color.rgb = VIOLET_700
    s.line.fill.background()

def hline(slide, left, top, width, color=VIOLET_700):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def card(slide, x, y, w, h, border_color=VIOLET_700):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = INK_900
    s.line.color.rgb = border_color
    s.line.width = Pt(1.5)
    return s

# ═══════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 1.4, 10, 1.2, 'ELIOS SCOGLIO', size=46, bold=True)
text(slide, 1.0, 2.5, 11, 0.7, 'Software & Architecture Manager  |  Tech Lead', size=22, color=VIOLET_400)
hline(slide, 1.0, 3.35, 3.2)
text(slide, 1.0, 3.65, 11, 0.55, 'Candidatura per il ruolo Tech Lead — Veralab / Re-Forme SRL', size=17, color=INK_200)
text(slide, 1.0, 4.25, 11, 0.45, 'Ho analizzato il vostro stack prima di arrivare qui.', size=16, color=INK_200, italic=True)
text(slide, 1.0, 6.2, 11, 0.5, '108 Vision — Costruiamo la direzione, non solo il codice.', size=13, color=INK_800)

# ═══════════════════════════════════════════════════
# SLIDE 2 — COSA HO TROVATO
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Cosa ho trovato (prima che me lo raccontiate)', size=34, bold=True)
bar(slide, 1.3)

# Left: funziona gia
text(slide, 1.2, 1.7, 5.3, 0.4, 'GIA PRESENTE', size=12, bold=True, color=GREEN_400)
works = [
    'Shopify — piattaforma solida e scalabile',
    'Jebbit skin diagnostic — personalizzazione CX',
    'Virtual try-on Overskin',
    'VERABILIA loyalty — meccanica definita',
    '14 store con Beauty Expert qualificati',
    'Videoconsulenze gratuite Face + Corpo',
    'Magazine SEO aggiornato settimanalmente',
]
tf = text(slide, 1.2, 2.15, 5.3, 4.0, '', size=14, color=INK_200)
for item in works:
    add_para(tf, f'✓  {item}', size=14, color=GREEN_400, space_before=7)

# Right: manca ancora
text(slide, 7.2, 1.7, 5.5, 0.4, 'DA COSTRUIRE', size=12, bold=True, color=AMBER_400)
gaps = [
    'Click & collect — non implementato',
    'Ship-from-store — non implementato',
    'Stock unificato online + 14 store',
    'Loyalty cross-channel (in-store)',
    'CDP / CRM unificato multi-touchpoint',
    'DevOps e CI/CD strutturati',
    'Monitoring & alert (500 errors silenti)',
]
tf2 = text(slide, 7.2, 2.15, 5.5, 4.0, '', size=14, color=INK_200)
for item in gaps:
    add_para(tf2, f'→  {item}', size=14, color=AMBER_400, space_before=7)

text(slide, 1.0, 6.6, 11, 0.4,
     'Ho trovato HTTP 500 su endpoint Shopify standard durante la preparazione. Probabilmente non avete alert su questo.',
     size=12, color=RED_400)

# ═══════════════════════════════════════════════════
# SLIDE 3 — CHI SONO
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Chi sono', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.7, 11, 0.45,
     'Software & Architecture Manager — piattaforme mission-critical da 10+ anni.',
     size=16, color=INK_200)

# Left: experience items
exp = [
    ('30M transazioni/anno', 'zero-downtime, spike gestiti, inventory real-time'),
    ('93 componenti, 7 livelli', '3 team coordinati sulla stessa roadmap'),
    ('Legacy 23 anni modernizzato', 'CORBA → microservizi gRPC — senza fermarlo'),
    ('Integration complesse', '7+ sistemi esterni: enti statali, fiscale, GDPR'),
]
for i, (title, desc) in enumerate(exp):
    y = 2.3 + i * 0.88
    text(slide, 1.5, y, 4.2, 0.4, title, size=16, bold=True, color=VIOLET_400)
    text(slide, 5.9, y, 6.8, 0.4, desc, size=15, color=INK_200)

# Results strip
hline(slide, 1.2, 5.95, 11.0)
results = [
    ('Deploy freq.', '+400%'),
    ('Durata deploy', '-91%'),
    ('Analisi requisiti', '-92%'),
    ('Team satisfaction', '+50%'),
]
for i, (label, number) in enumerate(results):
    x = 1.4 + i * 3.0
    text(slide, x, 6.1, 2.5, 0.35, label, size=12, color=INK_800)
    text(slide, x, 6.5, 2.5, 0.5, number, size=22, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════
# SLIDE 4 — IL GAP CENTRALE: INTEGRATION
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Il gap piu urgente: l\'integrazione', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.7, 11, 0.4,
     'La sfida non e Shopify — Shopify funziona. La sfida e che Shopify non parla con Dynamics NAV.',
     size=16, color=INK_200)

# Oggi box
card(slide, 1.0, 2.25, 5.3, 4.0, RED_400)
text(slide, 1.3, 2.4, 4.8, 0.4, 'OGGI (probabile)', size=14, bold=True, color=RED_400)
tf = text(slide, 1.3, 2.85, 4.8, 3.0, '', size=13, color=INK_200)
for line in [
    'Shopify  →  stock online',
    'NAV       →  stock fisico (14 store)',
    '',
    'Sync: notturno o manuale',
    '',
    'Risultato:',
    '  ✗  Nessuno sa dove sta lo stock real-time',
    '  ✗  Click & collect impossibile',
    '  ✗  Ship-from-store impossibile',
    '  ✗  Loyalty cross-channel impossibile',
]:
    add_para(tf, line, size=13, color=INK_200 if not line.startswith('  ✗') else RED_400, space_before=5)

# Target box
card(slide, 6.9, 2.25, 5.3, 4.0, GREEN_400)
text(slide, 7.2, 2.4, 4.8, 0.4, 'TARGET', size=14, bold=True, color=GREEN_400)
tf2 = text(slide, 7.2, 2.85, 4.8, 3.0, '', size=13, color=INK_200)
for line in [
    'Middleware Shopify ↔ NAV',
    'Inventory sync real-time',
    '',
    'Risultato:',
    '  ✓  Stock unificato (online + store + B2B)',
    '  ✓  Click & collect attivo',
    '  ✓  Ship-from-store attivo',
    '  ✓  Loyalty cross-channel',
    '  ✓  Profilo cliente unificato',
]:
    add_para(tf2, line, size=13, color=INK_200 if not line.startswith('  ✓') else GREEN_400, space_before=5)

text(slide, 1.0, 6.6, 11, 0.4,
     'Questo e il Problema #1. Tutto l\'omnichannel dipende da risolverlo prima.',
     size=13, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════
# SLIDE 5 — DECISIONE ERP
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'La decisione ERP da non rimandare', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.7, 11, 0.4,
     'Dynamics NAV e fuori mainstream support su versioni pre-2018. La decisione arriva comunque — meglio prenderla con dati.',
     size=15, color=INK_200)

options = [
    ('Migra a Business Central', 'API moderne, cloud-native, supporto MS', '6-18 mesi, migrazione dati', 'Urgente valutare', VIOLET_400),
    ('Resta su NAV + patch', 'Continuita, zero rischio migrazione', 'Lock-in crescente, API limitate', 'Costo nascosto crescente', AMBER_400),
    ('Nuovo ERP (NetSuite/SAP B1)', 'Liberta massima in futuro', 'Rischio altissimo a 70M revenue', 'Sconsigliato ora', RED_400),
]

headers = ['Opzione', 'Pro', 'Contro', 'Nota']
col_x = [1.2, 3.6, 6.5, 9.5]
col_w = [2.2, 2.7, 2.8, 2.5]
for j, h in enumerate(headers):
    text(slide, col_x[j], 2.35, col_w[j], 0.35, h, size=12, bold=True, color=INK_800)

hline(slide, 1.2, 2.75, 11.0, INK_800)

for i, (name, pro, con, note, col) in enumerate(options):
    y = 3.0 + i * 1.1
    text(slide, col_x[0], y, col_w[0], 0.85, name, size=14, bold=True, color=col)
    text(slide, col_x[1], y, col_w[1], 0.85, pro,  size=13, color=GREEN_400)
    text(slide, col_x[2], y, col_w[2], 0.85, con,  size=13, color=RED_400)
    text(slide, col_x[3], y, col_w[3], 0.85, note, size=13, color=INK_200)

text(slide, 1.0, 6.6, 11, 0.4,
     'Prima azione: Assessment NAV → Business Central nei primi 60 giorni. Non decidere a buio.',
     size=13, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════
# SLIDE 6 — ROADMAP 3 PRIORITA
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Il mio approccio: tre priorita in sequenza', size=34, bold=True)
bar(slide, 1.3)

phases = [
    ('0-90 gg', 'FOUNDATION', VIOLET_700, [
        'Audit stack: Shopify, NAV, CRM, WMS, POS',
        'Integration map — dove i dati si rompono',
        'Inventory sync Shopify ↔ NAV',
        'Monitoring + alert (stop ai 500 silenziosi)',
        'ADR per ogni decisione architetturale',
        'Assessment NAV → Business Central',
    ]),
    ('90-180 gg', 'OMNICHANNEL', VIOLET_500, [
        'Click & collect su tutti i 14 store',
        'Ship-from-store attivo',
        'Loyalty VERABILIA cross-channel',
        'Profilo cliente unificato (online + offline)',
        'CI/CD pipeline strutturata',
    ]),
    ('180-365 gg', 'SCALE', VIOLET_400, [
        'Multi-brand architecture (Veralab + Overskin)',
        'CDP: dati aggregati da tutti i touchpoint',
        'Migrazione ERP (se assessment lo conferma)',
        'AI use case su dati unificati',
        'Team building + governance stabile',
    ]),
]

for i, (period, title, col, items) in enumerate(phases):
    x = 0.6 + i * 4.2
    card(slide, x, 1.9, 3.9, 4.9, col)
    text(slide, x + 0.2, 2.05, 3.5, 0.35, period, size=11, bold=True, color=col)
    text(slide, x + 0.2, 2.45, 3.5, 0.5,  title,  size=19, bold=True)
    tf = text(slide, x + 0.2, 3.05, 3.5, 3.5, '', size=13, color=INK_200)
    for item in items:
        add_para(tf, f'→ {item}', size=13, color=INK_200, space_before=7)

text(slide, 1.0, 6.85, 11, 0.35,
     'Prima misuro, poi prometto. I numeri vengono dalla baseline — non da una slide.',
     size=12, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 7 — MATCH
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Match: il vostro contesto, la mia esperienza', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.5, 1.75, 5.0, 0.35, 'VERALAB', size=12, bold=True, color=INK_800)
text(slide, 7.8, 1.75, 5.0, 0.35, 'ELIOS SCOGLIO', size=12, bold=True, color=VIOLET_400)
hline(slide, 1.2, 2.15, 11.0, INK_800)

matches = [
    ('Shopify + NAV + 14 store + 870 touchpoint', 'Integration architecture, middleware, API-first'),
    ('Picchi traffico su lanci prodotto', 'On-sale ad alto volume: spike, inventory contesa, checkout sotto pressione'),
    ('Inventory non unificato', 'Real-time inventory sync su sistemi distribuiti'),
    ('Multi-brand Veralab + Overskin', 'Multi-tenant con shared infra e domain separation'),
    ('Compliance GDPR', 'GDPR quotidiano: PII, audit trail, data minimization'),
    ('Team lean, scope enorme', 'Cognitive load management, delivery cadenzata, ADR'),
    ('Nessun monitoring strutturato', 'Golden Signals, SLO, alert su ogni servizio critico'),
]

for i, (their, mine) in enumerate(matches):
    y = 2.3 + i * 0.62
    text(slide, 1.5, y, 5.0, 0.5, their, size=13, color=INK_200)
    hline(slide, 6.6, y + 0.2, 0.9)
    text(slide, 7.8, y, 5.0, 0.5, mine,  size=13, bold=True, color=VIOLET_400)

# ═══════════════════════════════════════════════════
# SLIDE 8 — OPPORTUNITA AI
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'L\'opportunita che la JD non cita: AI', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.7, 11, 0.4,
     'Il brand ha gia la base per fare AI in modo serio — manca solo la foundation dati unificata.',
     size=15, color=INK_200)

# Data assets
text(slide, 1.2, 2.25, 5.5, 0.35, 'ASSET DATI GIA PRESENTI', size=12, bold=True, color=VIOLET_400)
assets = [
    'Skin diagnostic (Jebbit) = profilo cliente strutturato',
    'VERABILIA = storico acquisti + comportamento fedele',
    'Magazine = corpus contenuti per RAG',
    '14 store = dati comportamentali offline',
    '870 touchpoint = segnali di acquisto B2B',
]
tf = text(slide, 1.2, 2.65, 5.5, 2.8, '', size=14, color=INK_200)
for item in assets:
    add_para(tf, f'→  {item}', size=13, color=GREEN_400, space_before=7)

# AI use cases
text(slide, 7.2, 2.25, 5.7, 0.35, 'USE CASE AD ALTO ROI', size=12, bold=True, color=VIOLET_400)
cases = [
    ('Recommendation personalizzata', 'skin type + acquisti → +conversion +AOV'),
    ('Chatbot beauty expert 24/7', 'RAG sulla KB prodotti → scalare le consulenze'),
    ('Segmentazione predittiva', 'propensity Veralab→Overskin → revenue incrementale'),
]
for i, (title, desc) in enumerate(cases):
    y = 2.65 + i * 1.05
    card(slide, 7.2, y, 5.7, 0.9, VIOLET_700)
    text(slide, 7.4, y + 0.08, 5.3, 0.35, title, size=14, bold=True, color=VIOLET_400)
    text(slide, 7.4, y + 0.48, 5.3, 0.32, desc, size=12, color=INK_200)

text(slide, 1.0, 6.05, 11, 0.7,
     '"Il prerequisito di ogni use case AI e uno solo: dati unificati.\nCostruire prima la foundation — poi l\'AI funziona davvero."',
     size=15, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 9 — METODO
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Il metodo: Misura, Decide, Esegui', size=34, bold=True)
bar(slide, 1.3)

steps = [
    ('MESE 1', 'ASCOLTO', VIOLET_700, [
        'Non cambio niente in produzione.',
        'Mappo lo stack completo (Shopify, NAV, CRM, WMS, POS, loyalty).',
        'Trovo dove i dati si rompono tra sistemi.',
        'Output: Tech Map + Roadmap 12 mesi prioritizzata.',
    ]),
    ('MESI 2-3', 'PRIME AZIONI', VIOLET_500, [
        'Inventory sync Shopify ↔ NAV (primo mattone).',
        'Monitoring + alert di base su Shopify.',
        'ADR per le prime 5 decisioni architetturali.',
        'Assessment NAV → Business Central.',
    ]),
    ('MESI 4-6', 'OMNICHANNEL', VIOLET_400, [
        'Click & collect attivo su tutti i 14 store.',
        'Loyalty cross-channel VERABILIA.',
        'Profilo cliente unificato online + offline.',
        'CI/CD pipeline strutturata.',
    ]),
]

for i, (period, title, col, items) in enumerate(steps):
    x = 0.6 + i * 4.2
    card(slide, x, 1.85, 3.9, 4.4, col)
    text(slide, x + 0.2, 2.0,  3.5, 0.35, period, size=11, bold=True, color=col)
    text(slide, x + 0.2, 2.38, 3.5, 0.5,  title,  size=19, bold=True)
    tf = text(slide, x + 0.2, 2.95, 3.5, 3.0, '', size=13, color=INK_200)
    for item in items:
        add_para(tf, item, size=13, color=INK_200, space_before=9)

text(slide, 1.0, 6.7, 11, 0.4,
     'Prima misuro, poi prometto. Mai numeri senza baseline reale.',
     size=13, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 10 — UNA COSA CONCRETA (500 errors)
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Una cosa concreta — trovata prima di questo colloquio', size=34, bold=True)
bar(slide, 1.3)

card(slide, 1.5, 1.85, 10.0, 3.2, RED_400)
text(slide, 2.0, 2.05, 9.0, 0.45, 'PROBLEMA RILEVATO: HTTP 500 su Shopify', size=18, bold=True, color=RED_400)
tf = text(slide, 2.0, 2.6, 9.0, 2.2, '', size=15, color=INK_200)
for line in [
    'Endpoint coinvolti:  /blogs/news  —  API standard Shopify',
    '',
    'Causa probabile: conflitto tra app di terze parti o liquid template',
    'Impatto: contenuti magazine non raggiungibili, API instabili',
    'Monitoring attuale: probabilmente nessuno (nessun alert visibile)',
]:
    add_para(tf, line, size=15, color=INK_200 if not line.startswith('Endpoint') else RED_400, space_before=8)

text(slide, 1.2, 5.3, 11, 0.5,
     'Non e un\'emergenza. E un segnale: la manutenzione proattiva dello stack non e presidiata.',
     size=16, color=AMBER_400)

text(slide, 1.2, 5.95, 11, 0.75,
     'Il primo atto da Tech Lead e sistemarla — non perche i 500 siano critici,\nma perche avere visibilita completa e il prerequisito di tutto il resto.',
     size=15, color=INK_200)

# ═══════════════════════════════════════════════════
# SLIDE 11 — TRE DOMANDE
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Tre domande che voglio fare a voi', size=34, bold=True)
bar(slide, 1.3)

text(slide, 1.2, 1.75, 11, 0.4,
     'Le risposte cambiano la priorita della roadmap. Le voglio capire oggi, non in un report.',
     size=15, color=INK_200, italic=True)

questions = [
    (
        '1',
        'Come funziona oggi la sincronizzazione inventory tra Shopify e i 14 store?',
        'Questa risposta determina quanto e urgente e complessa la priority #1.',
        VIOLET_700,
    ),
    (
        '2',
        'Overskin e sullo stesso account Shopify o e uno store separato?',
        'Cambia radicalmente l\'architettura multi-brand e il data model del cliente.',
        VIOLET_500,
    ),
    (
        '3',
        'Chi ha preso le decisioni architetturali finora? Esiste un backlog tecnico?',
        'Determina il punto di partenza: costruire da zero o evolvere una direzione esistente.',
        VIOLET_400,
    ),
]

for i, (num, q, why, col) in enumerate(questions):
    y = 2.55 + i * 1.45
    card(slide, 1.2, y, 10.9, 1.25, col)
    text(slide, 1.55, y + 0.1,  0.6, 0.5, num, size=22, bold=True, color=col)
    text(slide, 2.3,  y + 0.1, 9.4, 0.5, q,   size=16, bold=True)
    text(slide, 2.3,  y + 0.62, 9.4, 0.45, why, size=13, color=INK_800)

# ═══════════════════════════════════════════════════
# SLIDE 12 — NEXT STEP + CHIUSURA
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Prossimo passo', size=34, bold=True)
bar(slide, 1.3)

card(slide, 1.5, 1.85, 10.0, 3.0, VIOLET_700)
text(slide, 2.0, 2.05, 9.0, 0.5, 'MESE 1 — TECH MAP', size=22, bold=True, color=VIOLET_400)
tf = text(slide, 2.0, 2.65, 9.0, 1.9, '', size=15, color=INK_200)
for line in [
    'Durata: 30 giorni — nessuna modifica in produzione',
    'Output: mappa completa stack, integration map, top 5 rischi, roadmap 12 mesi',
    'Poi: obiettivi concordati con KPI misurabili, review a fine trimestre',
]:
    add_para(tf, f'→  {line}', size=15, color=INK_200, space_before=9)

text(slide, 1.0, 5.2, 11, 0.7,
     '"Non vengo a portare teoria.\nVengo con un\'analisi concreta del vostro stack e tre problemi che posso risolvere."',
     size=18, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

text(slide, 1.0, 6.0, 11, 0.45,
     '"Il primo mese ascolto. Dal secondo, i numeri parlano."',
     size=16, color=INK_200, align=PP_ALIGN.CENTER, italic=True)

text(slide, 1.0, 6.65, 11, 0.4,
     'Elios Scoglio  |  elios@108vision.it  |  108vision.it',
     size=13, color=INK_800, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 13 — SICUREZZA / PERFORMANCE / RESILIENZA
# ═══════════════════════════════════════════════════
slide = dark_slide()
text(slide, 1.0, 0.5, 11, 0.75, 'Analisi tecnica: Sicurezza, Performance, Resilienza', size=32, bold=True)
bar(slide, 1.3)

# ---- colonna sinistra: Sicurezza ----
text(slide, 1.2, 1.75, 3.7, 0.35, 'SECURITY HEADERS', size=11, bold=True, color=VIOLET_400)
sec_items = [
    (GREEN_400,  '✓  HSTS — Shopify impone HTTPS'),
    (GREEN_400,  '✓  X-Frame-Options — default Shopify'),
    (RED_400,    '✗  CSP — assente / permissivo'),
    (AMBER_400,  '~  Referrer-Policy — non configurato'),
    (RED_400,    '✗  Permissions-Policy — assente'),
    (AMBER_400,  '~  Cookie consent — non verificabile'),
]
tf = text(slide, 1.2, 2.15, 3.7, 3.2, '', size=13, color=INK_200)
for col_s, item in sec_items:
    add_para(tf, item, size=13, color=col_s, space_before=8)

text(slide, 1.2, 5.45, 3.7, 0.5,
     'Jebbit (USA) = dati skin type\nSCCs GDPR obbligatori',
     size=11, color=AMBER_400)

# ---- colonna centrale: Performance ----
text(slide, 5.3, 1.75, 3.6, 0.35, 'CORE WEB VITALS', size=11, bold=True, color=VIOLET_400)
cwv = [
    ('LCP',  '2.5 – 4s',   AMBER_400, 'video hero senza priorita'),
    ('CLS',  '< 0.1',      GREEN_400, 'immagini con dimensioni dichiarate'),
    ('INP',  '200 – 500ms',AMBER_400, 'bundle JS 3rd party pesante'),
    ('TTFB', '< 200ms',    GREEN_400, 'Shopify CDN EU performante'),
]
for i, (metric, val, col_c, note) in enumerate(cwv):
    y = 2.15 + i * 0.85
    text(slide, 5.3, y,      1.0, 0.4, metric, size=14, bold=True, color=col_c)
    text(slide, 6.4, y,      2.0, 0.4, val,    size=14, bold=True, color=col_c)
    text(slide, 5.3, y+0.42, 3.6, 0.3, note,   size=11, color=INK_800)

card(slide, 5.1, 5.0, 3.9, 0.95, RED_400)
text(slide, 5.3, 5.1, 3.6, 0.35, 'HTTP 500 RILEVATI LIVE', size=13, bold=True, color=RED_400)
text(slide, 5.3, 5.5, 3.6, 0.35, '/blogs/news  •  /products.json  •  /cart.js', size=11, color=INK_200)

# ---- colonna destra: Resilienza ----
text(slide, 9.4, 1.75, 3.6, 0.35, 'RESILIENZA / SPOF', size=11, bold=True, color=VIOLET_400)
spofs = [
    (RED_400,    '✗  NAV offline = gestione ordini bloccata'),
    (RED_400,    '✗  Nessun circuit breaker su Shopify↔NAV'),
    (AMBER_400,  '~  Shopify 99.98% SLA (vendor lock-in)'),
    (AMBER_400,  '~  API rate limit (40 req/s) durante picchi'),
    (GREEN_400,  '✓  Shopify CDN — edge PoP globali'),
    (GREEN_400,  '✓  robots.txt protegge checkout da crawl'),
]
tf3 = text(slide, 9.4, 2.15, 3.7, 3.2, '', size=12, color=INK_200)
for col_r, item in spofs:
    add_para(tf3, item, size=12, color=col_r, space_before=8)

text(slide, 9.4, 5.45, 3.7, 0.5,
     'Peak traffic = rischio desync stock\nPattern identico all\'on-sale ticketing',
     size=11, color=AMBER_400)

# ---- footer priorita ----
hline(slide, 1.0, 6.05, 11.2, INK_800)
text(slide, 1.0, 6.2, 11, 0.5,
     'P0: fix HTTP 500 + circuit breaker su NAV  |  P1: CSP + GDPR Jebbit  |  P2: bundle JS audit  |  P3: structured data',
     size=12, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
out_path = r'c:\Code\Documents\Lavoro\Personale\Vision\tracks\108-cto\PRES_Veralab_TechLead.pptx'
prs.save(out_path)
print(f'PowerPoint salvato: {out_path}')
