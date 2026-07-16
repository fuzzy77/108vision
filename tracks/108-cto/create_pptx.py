from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 108 Vision palette
INK_950 = RGBColor(0x0F, 0x17, 0x2A)
INK_900 = RGBColor(0x1E, 0x29, 0x3B)
INK_800 = RGBColor(0x33, 0x41, 0x55)
INK_200 = RGBColor(0xE2, 0xE8, 0xF0)
INK_50 = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
VIOLET_700 = RGBColor(0x6D, 0x28, 0xD9)
VIOLET_500 = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_400 = RGBColor(0xA7, 0x8B, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = INK_950
    return slide

def text(slide, left, top, width, height, txt, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = 'Inter'
    return tf

def bar(slide, top):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(0.08), Inches(0.8))
    s.fill.solid()
    s.fill.fore_color.rgb = VIOLET_700
    s.line.fill.background()

def hline(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = VIOLET_700
    s.line.fill.background()

# ===== SLIDE 1: COVER =====
slide = dark_slide()
text(slide, 1.0, 1.5, 10, 1.2, 'ELIOS SCOGLIO', size=44, bold=True)
text(slide, 1.0, 2.5, 10, 0.8, 'Software & Architecture Manager | Fractional CTO', size=22, color=VIOLET_400)
hline(slide, 1.0, 3.5, 3.0)
text(slide, 1.0, 4.0, 10, 0.6, 'Presentazione per Etica Soluzioni', size=18, color=INK_200)
text(slide, 1.0, 6.2, 10, 0.5, '108 Vision — Costruiamo la direzione, non solo il codice.', size=14, color=INK_800)

# ===== SLIDE 2: CHI E' 108 VISION =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Chi e 108 Vision', size=36, bold=True)
bar(slide, 1.4)
text(slide, 1.2, 1.8, 10, 0.5, '25 anni di esperienza su sistemi in produzione. Risultati misurati, non promessi.', size=16, color=INK_200)

# Left column: experiences
text(slide, 1.2, 2.5, 5.5, 0.4, 'COSA ABBIAMO FATTO', size=13, bold=True, color=VIOLET_400)
experiences = [
    'Piattaforma immobiliare nazionale — refactoring completo processi di delivery',
    'Cloud provider enterprise — trasformazione qualita e rilascio su larga scala',
    'Piattaforma ticketing 30M transazioni/anno — governance architetturale, 93 componenti, 3 team',
    'Compliance mission-critical — fiscale, Polizia di Stato, GDPR, PCI',
    'Legacy modernization — sistema 23 anni (CORBA → microservizi gRPC)',
    'AI adoption pragmatica — integrata nel ciclo di sviluppo con ROI dimostrato',
]
tf = text(slide, 1.2, 3.0, 5.8, 3.5, '', size=14, color=INK_200)
for item in experiences:
    p = tf.add_paragraph()
    p.text = f'→  {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(8)

# Right column: results
text(slide, 7.5, 2.5, 5.0, 0.4, 'RISULTATI MISURATI', size=13, bold=True, color=VIOLET_400)
results_list = [
    ('Lead time', '-60%'),
    ('Bug rate', '-98%'),
    ('Deploy frequency', '+400%'),
    ('Tempo deploy', '-91%'),
    ('Costo sviluppo (AI)', '-77%'),
    ('Team satisfaction', '+50%'),
]
for i, (label, number) in enumerate(results_list):
    y = 3.0 + i * 0.6
    text(slide, 7.5, y, 3.0, 0.4, label, size=14, color=INK_200)
    text(slide, 10.5, y, 2.0, 0.4, number, size=18, bold=True, color=VIOLET_400)

text(slide, 1.0, 6.6, 11, 0.4, 'Costruiamo la direzione, non solo il codice. Il valore e nelle decisioni giuste, non nelle righe scritte.', size=14, bold=True, color=VIOLET_400)

# ===== SLIDE 3: COSA FACCIO =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Cosa faccio concretamente', size=36, bold=True)
bar(slide, 1.4)
areas = [
    ('Strategia', 'Roadmap tecnica, decisioni architetturali, allineamento business-tech'),
    ('Architettura', 'Review design, standard, debito tecnico, scelte build/buy'),
    ('Team', 'Hiring, 1:1, mentoring, cultura ingegneristica, crescita'),
    ('Governance', 'ADR tracciabili, metriche DORA, report mensile, CEO sync'),
    ('AI & Innovazione', 'Adoption pragmatica — ROI dimostrabile, non hype'),
]
for i, (title, desc) in enumerate(areas):
    y = 2.0 + i * 0.95
    text(slide, 1.5, y, 3.0, 0.5, title, size=18, bold=True, color=VIOLET_400)
    text(slide, 5.0, y, 7.5, 0.5, desc, size=16, color=INK_200)
text(slide, 1.0, 6.6, 10, 0.4, 'Cosa NON faccio: scrivere codice, gestire sprint, fare il PM.', size=13, color=INK_800)

# ===== SLIDE 4: RISULTATI =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Risultati dimostrabili', size=36, bold=True)
bar(slide, 1.4)
results = [
    ('Deploy frequency', '+400%', 'da ogni 6 settimane a ogni settimana'),
    ('Tempo deploy', '-91%', 'da 4 ore a 22 minuti'),
    ('Costo sviluppo (AI)', '-77/82%', 'su task specifici'),
    ('Tempo analisi', '-92%', 'da 2-3 giorni a 2 ore'),
    ('Team satisfaction', '+50%', 'survey interna'),
]
for i, (label, number, detail) in enumerate(results):
    y = 2.2 + i * 0.95
    text(slide, 1.5, y, 3.5, 0.5, label, size=16, color=INK_200)
    text(slide, 5.2, y, 2.0, 0.5, number, size=24, bold=True, color=VIOLET_400)
    text(slide, 7.5, y, 5.0, 0.5, detail, size=14, color=INK_800)
text(slide, 1.0, 6.6, 10, 0.4, 'Numeri da sistemi reali in produzione. Non POC, non demo.', size=13, color=INK_800)

# ===== SLIDE 5: DUE FILONI =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Il mio approccio: due filoni paralleli', size=36, bold=True)
bar(slide, 1.4)
# Filone 1
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(5.3), Inches(3.8))
s.fill.solid()
s.fill.fore_color.rgb = INK_900
s.line.color.rgb = VIOLET_700
s.line.width = Pt(1.5)
text(slide, 1.3, 2.4, 4.8, 0.5, 'FILONE 1 — PRESENTE', size=16, bold=True, color=VIOLET_400)
text(slide, 1.3, 2.9, 4.8, 0.4, 'Concretizzare e ottimizzare', size=14, color=INK_200)
tf = text(slide, 1.3, 3.4, 4.8, 2.5, '', size=14, color=INK_200)
for item in ['Flussi di sviluppo ottimizzati', 'Best practice e standard condivisi', 'CI/CD e testing automatico', 'Incident management strutturato', 'Comunicazione team-business']:
    p = tf.add_paragraph()
    p.text = f'→ {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(6)
text(slide, 1.3, 5.6, 4.8, 0.4, 'Valore visibile in 2-4 settimane.', size=12, bold=True, color=VIOLET_400)

# Filone 2
s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.3), Inches(3.8))
s2.fill.solid()
s2.fill.fore_color.rgb = INK_900
s2.line.color.rgb = VIOLET_500
s2.line.width = Pt(1.5)
text(slide, 7.1, 2.4, 4.8, 0.5, 'FILONE 2 — VISIONE', size=16, bold=True, color=VIOLET_400)
text(slide, 7.1, 2.9, 4.8, 0.4, 'Dove andare nei prossimi 3-5 anni', size=14, color=INK_200)
tf2 = text(slide, 7.1, 3.4, 4.8, 2.5, '', size=14, color=INK_200)
for item in ['Revisione architetturale completa', 'AI strategy (prodotto + processo)', 'Scale engineering (2x clienti)', 'Cloud & platform approach', 'Security operativa (code-level)']:
    p = tf2.add_paragraph()
    p.text = f'→ {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(6)
text(slide, 7.1, 5.6, 4.8, 0.4, 'Decisioni misurabili con ADR.', size=12, bold=True, color=VIOLET_400)

# ===== SLIDE 5B: PRINCIPI — OVERVIEW =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'I principi dietro ogni decisione', size=36, bold=True)
bar(slide, 1.4)
text(slide, 1.2, 2.0, 10, 0.6, '15 pilastri tecnici organizzati in 4 macro-aree:', size=16, color=INK_200)
pillars = [
    ('COSTRUIRE', 'Architettura, Design, Qualita', 'Cosa stiamo costruendo e come?'),
    ('VERIFICARE', 'Testing, Delivery, Osservabilita', 'Come sappiamo che funziona?'),
    ('PROTEGGERE', 'Sicurezza, Dati, Scalabilita', 'Come lo difendiamo?'),
    ('GUIDARE', 'Team, Debito, AI, Decisioni, Costi, Comunicazione', 'Come funzionano le persone?'),
]
for i, (area, topics, question) in enumerate(pillars):
    y = 2.8 + i * 1.05
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(y), Inches(10.8), Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = INK_900
    box.line.color.rgb = VIOLET_700
    box.line.width = Pt(1)
    text(slide, 1.5, y + 0.08, 2.8, 0.4, area, size=16, bold=True, color=VIOLET_400)
    text(slide, 4.5, y + 0.08, 4.5, 0.4, topics, size=14, color=INK_200)
    text(slide, 9.2, y + 0.08, 2.8, 0.4, question, size=13, color=INK_800)
text(slide, 1.0, 6.7, 11, 0.4, 'Non teoria accademica. Framework operativo testato su sistemi enterprise reali.', size=13, bold=True, color=VIOLET_400)

# ===== SLIDE 5C: PRINCIPI — COSTRUIRE =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'COSTRUIRE — Principi architetturali', size=36, bold=True)
bar(slide, 1.4)
build_principles = [
    ('Trade-off espliciti', 'Non esiste giusto. Esiste adatto a un costo dichiarato.'),
    ('ADR obbligatori', 'Se non e scritto, tra 6 mesi diventa folklore.'),
    ('Bounded Context', 'Progetta dai linguaggi del business, non dalle tabelle.'),
    ('Resilienza B-T-R-C', 'Bulkhead → Timeout → Retry → Circuit Breaker. Sempre.'),
    ('API-First', 'Contratto prima, codice dopo. Il sistema scala sul contratto.'),
    ('Modular Monolith', 'Non estrarre un servizio finche il confine non esiste gia.'),
]
for i, (title, desc) in enumerate(build_principles):
    y = 2.0 + i * 0.82
    text(slide, 1.5, y, 4.0, 0.4, title, size=16, bold=True, color=VIOLET_400)
    text(slide, 5.8, y, 6.5, 0.4, desc, size=15, color=INK_200)
text(slide, 1.0, 6.7, 11, 0.4, 'Ogni principio ha una fitness function misurabile e un ADR di supporto.', size=13, color=INK_800)

# ===== SLIDE 5D: PRINCIPI — VERIFICARE =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'VERIFICARE — Dal codice alla produzione', size=36, bold=True)
bar(slide, 1.4)
verify_principles = [
    ('Testing', 'Coverage alto ≠ fiducia alta.\nTesta il comportamento, non come e scritto dentro.'),
    ('Delivery', 'Deploy e un evento tecnico noioso.\nRelease e una decisione di prodotto.'),
    ('Osservabilita', 'In produzione non hai il debugger.\nHai solo Log, Metriche, Trace.'),
]
for i, (title, desc) in enumerate(verify_principles):
    y = 2.2 + i * 1.5
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(y), Inches(10.8), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = INK_900
    box.line.color.rgb = VIOLET_700
    box.line.width = Pt(1)
    text(slide, 1.6, y + 0.15, 3.0, 0.5, title, size=18, bold=True, color=VIOLET_400)
    text(slide, 5.0, y + 0.15, 6.8, 1.0, desc, size=15, color=INK_200)
verify_kpis = [
    'Pipeline < 10 min',
    'Rollback < 5 min',
    'Golden Signals su ogni servizio',
    'Contract test tra team',
]
tf = text(slide, 1.2, 6.0, 11, 0.6, '', size=13, color=INK_800)
for kpi in verify_kpis:
    p = tf.add_paragraph()
    p.text = f'✓  {kpi}'
    p.font.size = Pt(13)
    p.font.color.rgb = INK_800
    p.font.name = 'Inter'

# ===== SLIDE 5E: PRINCIPI — PROTEGGERE & GUIDARE =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'PROTEGGERE & GUIDARE', size=36, bold=True)
bar(slide, 1.4)
# Left column: Proteggere
text(slide, 1.2, 2.0, 5.0, 0.4, 'PROTEGGERE', size=14, bold=True, color=VIOLET_400)
protect_items = [
    'Security by Design — default = negare accesso',
    'SAST + dep scan + secret detection in CI',
    'PostgreSQL a meno che... (motivo specifico)',
    'Scala verticale finche costa meno',
    'Cache il 20% hot, load test a 2x picco',
]
tf = text(slide, 1.2, 2.5, 5.3, 3.5, '', size=14, color=INK_200)
for item in protect_items:
    p = tf.add_paragraph()
    p.text = f'→  {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(10)
# Right column: Guidare
text(slide, 7.0, 2.0, 5.5, 0.4, 'GUIDARE', size=14, bold=True, color=VIOLET_400)
guide_items = [
    'Cognitive load sostenibile (semplifica)',
    '15% sprint per debito — non negoziabile',
    'AI come sistema, non come demo',
    'Reversibile = decidi ora; Type 1 = ADR',
    'Mai sorprese, sempre opzioni, quantifica',
]
tf2 = text(slide, 7.0, 2.5, 5.5, 3.5, '', size=14, color=INK_200)
for item in guide_items:
    p = tf2.add_paragraph()
    p.text = f'→  {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(10)
text(slide, 1.0, 6.2, 11, 0.8, '"Non vendo teoria. Vendo un sistema operativo per decisioni tecniche\nche il CEO puo capire e il team puo eseguire."', size=15, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)

# ===== SLIDE 5F: PRINCIPI — VALORE PER VOI =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Cosa significa per Etica Soluzioni', size=36, bold=True)
bar(slide, 1.4)
value_mapping = [
    ('ADR + Fitness Functions', 'Decisioni architetturali tracciate e verificabili'),
    ('B-T-R-C su ogni integrazione PA', 'PagoPA/SPID/AppIO down non blocca il sistema'),
    ('Golden Signals + SLO', 'Sapete in 30 sec se c\'e un problema, non dal cliente'),
    ('15% debito + prioritizzazione', 'Il prodotto evolve senza accumulare costo nascosto'),
    ('Deploy ≠ Release + Feature Flags', 'Rilasciate quando volete, rollback in 5 minuti'),
    ('Cost routing AI', 'Il team AI produce risultati, non fatture sorpresa'),
]
for i, (principle, value) in enumerate(value_mapping):
    y = 2.0 + i * 0.78
    text(slide, 1.5, y, 5.0, 0.4, principle, size=15, bold=True, color=VIOLET_400)
    text(slide, 6.8, y, 5.5, 0.4, value, size=15, color=INK_200)
text(slide, 1.0, 6.7, 11, 0.4, 'Ogni principio si traduce in un risultato misurabile per il business.', size=13, bold=True, color=VIOLET_400)

# ===== SLIDE 6: METODO =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Il metodo: Know-How → Step → KPI', size=36, bold=True)
bar(slide, 1.4)
phases = [
    ('MESE 1', 'ASCOLTO', 'Capisco il vostro mondo.\nMappo architettura, flussi, team.\nMisuro la baseline.\nOutput: State of the Stack.'),
    ('MESI 2-3', 'AZIONE', 'Obiettivi concreti con owner\ne deadline. 3 step operativi\n+ 3 step strategici.\nKPI misurati mese per mese.'),
    ('MESE 4+', 'ESECUZIONE', 'Report mensile con numeri.\nRoadmap viva. Retrospettiva\nteam. CEO sync bisettimanale.\nScalare up/down a scelta.'),
]
for i, (period, title, desc) in enumerate(phases):
    x = 1.0 + i * 4.0
    colors = [VIOLET_700, VIOLET_500, VIOLET_400]
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(3.6), Inches(4.2))
    box.fill.solid()
    box.fill.fore_color.rgb = INK_900
    box.line.color.rgb = colors[i]
    box.line.width = Pt(2)
    text(slide, x + 0.2, 2.4, 3.2, 0.4, period, size=12, bold=True, color=colors[i])
    text(slide, x + 0.2, 2.9, 3.2, 0.5, title, size=22, bold=True)
    text(slide, x + 0.2, 3.6, 3.2, 2.5, desc, size=14, color=INK_200)
text(slide, 1.0, 6.7, 11, 0.4, 'Prima misuro, poi prometto. Mai numeri senza baseline.', size=13, bold=True, color=VIOLET_400)

# ===== SLIDE 7: FRACTIONAL vs FULL-TIME =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Fractional CTO vs Full-Time', size=36, bold=True)
bar(slide, 1.4)
comparisons = [
    ('Costo annuo', '120-180K EUR', '85-100K EUR'),
    ('Rischio hiring', '6 mesi per capire', 'Trial 3 mesi, zero rischio'),
    ('Prospettiva', 'Solo interna', 'Esterna + interna'),
    ('Competenza AI', 'Da trovare + formare', 'Disponibile dal giorno 1'),
    ('Flessibilita', 'Fisso (rinegoziare)', 'Up/down ogni mese'),
    ('Exit', 'Panico se esce', 'Pianificata = successo'),
]
text(slide, 4.5, 1.9, 3.8, 0.4, 'Full-time', size=13, bold=True, color=INK_800)
text(slide, 8.8, 1.9, 3.8, 0.4, 'Fractional', size=13, bold=True, color=VIOLET_400)
for i, (label, ft, frac) in enumerate(comparisons):
    y = 2.4 + i * 0.75
    text(slide, 1.2, y, 3.0, 0.4, label, size=15, bold=True)
    text(slide, 4.5, y, 3.8, 0.4, ft, size=14, color=INK_200)
    text(slide, 8.8, y, 3.8, 0.4, frac, size=14, bold=True, color=VIOLET_400)

# ===== SLIDE 8: MATCH =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Match con il vostro contesto', size=36, bold=True)
bar(slide, 1.4)
matches = [
    ('750K prenotazioni/giorno', '30M transazioni/anno'),
    ('1.000+ Comuni multi-tenant', 'Multi-tenant enterprise'),
    ('PagoPA, SPID, AppIO', 'SIAE, VRO, GDPR compliance'),
    ('Prodotto 23 anni', 'Legacy 23 anni in migrazione'),
    ('7+ integrazioni PA', '7+ integrazioni + circuit breaker'),
    ('Dati minori e pazienti', 'GDPR quotidiano, PII controls'),
    ('Team AI senza direzione', 'AI con ROI -77% dimostrato'),
]
text(slide, 1.5, 1.9, 5.0, 0.4, 'ETICA SOLUZIONI', size=12, bold=True, color=INK_800)
text(slide, 7.8, 1.9, 5.0, 0.4, 'ELIOS SCOGLIO', size=12, bold=True, color=VIOLET_400)
for i, (their, mine) in enumerate(matches):
    y = 2.4 + i * 0.68
    text(slide, 1.5, y, 5.0, 0.4, their, size=14, color=INK_200)
    text(slide, 7.8, y, 5.0, 0.4, mine, size=14, bold=True, color=VIOLET_400)
    hline(slide, 6.5, y + 0.2, 1.0)

# ===== SLIDE 9: COME FUNZIONA =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Come funziona in pratica', size=36, bold=True)
bar(slide, 1.4)
details = [
    ('Presenza', '2 giorni/settimana dedicati (scalabile)'),
    ('Giorno 1', 'Governance + Team: standup, 1:1, review architetturali'),
    ('Giorno 2', 'Strategia + Stakeholder: CEO sync, roadmap, metriche'),
    ('Async', 'Risposta entro 4h; urgenze entro 2h'),
    ('Ogni mese', 'Report + ADR + roadmap + metriche'),
    ('Commitment', '3 mesi minimo, poi mensile'),
    ('Entry point', 'Tech Assessment 2gg — zero vincolo'),
]
for i, (label, value) in enumerate(details):
    y = 2.0 + i * 0.72
    text(slide, 1.5, y, 3.0, 0.4, label, size=15, bold=True, color=VIOLET_400)
    text(slide, 4.8, y, 8.0, 0.4, value, size=15, color=INK_200)

# ===== SLIDE 10: NEXT STEP =====
slide = dark_slide()
text(slide, 1.0, 0.6, 10, 0.8, 'Prossimo passo', size=36, bold=True)
bar(slide, 1.4)
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(10.0), Inches(3.2))
box.fill.solid()
box.fill.fore_color.rgb = INK_900
box.line.color.rgb = VIOLET_700
box.line.width = Pt(2)
text(slide, 2.0, 2.3, 9.0, 0.5, 'TECH ASSESSMENT', size=24, bold=True, color=VIOLET_400)
tf = text(slide, 2.0, 3.0, 9.0, 2.0, '', size=16, color=INK_200)
for line in ['2 giorni — Output: State of the Stack', 'Architettura, team, flussi, rischi, priorita', 'Costo detratto dal primo mese se si prosegue', 'Zero vincolo: deliverable di valore comunque', '', 'Poi: 3 mesi → obiettivi + KPI → review trimestrale']:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = INK_200
    p.font.name = 'Inter'
    p.space_before = Pt(8)

text(slide, 1.0, 5.8, 11, 0.6, '"Il rischio non e provare. Il rischio e continuare senza direzione tecnica."', size=20, bold=True, color=VIOLET_400, align=PP_ALIGN.CENTER)
text(slide, 1.0, 6.6, 11, 0.4, 'Elios Scoglio | elios@108vision.it | 108vision.it', size=13, color=INK_800, align=PP_ALIGN.CENTER)

prs.save(r'c:\Code\Documents\Lavoro\Personale\Vision\tracks\108-cto\PRES_EticaSoluzioni_EliosScoglio.pptx')
print('PowerPoint salvato con successo.')
